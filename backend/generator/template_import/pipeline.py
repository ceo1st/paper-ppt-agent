"""Sequential state machine orchestrating a template import.

Implements the ``uploaded → analyzing → rendering → detecting_assets →
llm_review → review`` flow described in ``design.md``. Each step
delegates to a leaf module and is wrapped by ``observability.step_timing``;
failures are mapped to the ``error_kind`` taxonomy
(``render`` / ``extraction`` / ``llm`` / ``persistence`` / ``unknown``).
"""

from __future__ import annotations

import base64
import hashlib
import html
import json
import logging
import mimetypes
import re
import time
import xml.etree.ElementTree as ET
from pathlib import Path, PurePosixPath
from typing import Any

from . import asset_extractor, chrome_detector, page_classifier, persistence, templateizer
from . import renderer as renderer_pkg
from .llm_client import LLMClient, merge_llm_plan
from .observability import step_timing
from .types import (
    AssetCandidate,
    BoundingBox,
    ChromeTextItem,
    ElementAction,
    ElementActionRecord,
    ExtractionError,
    ImportResult,
    ImportTaskState,
    LayoutPack,
    LLMPlanError,
    LLMTraceEntry,
    ManifestAssetEntry,
    ManifestCanvas,
    ManifestContentArea,
    ManifestTheme,
    ManifestWarning,
    PageType,
    PersistenceError,
    PipelineContext,
    RenderError,
    RenderResult,
    ReviewDraft,
    SlideMetrics,
    StepRecord,
    TemplateImportError,
    TemplateManifest,
    TemplateizeOutput,
)


logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Step table (id, label, default error_kind, skippable)
# ─────────────────────────────────────────────────────────────────────────────

_STEPS: tuple[tuple[str, str, str, bool], ...] = (
    ("uploaded",         "Upload received",     "extraction", False),
    ("analyzing",        "Analyzing PPTX",      "extraction", False),
    ("rendering",        "Rendering slides",    "render",     False),
    ("detecting_assets", "Detecting assets",    "extraction", False),
    ("llm_review",       "LLM review",          "llm",        True),
    ("review",           "Awaiting review",     "unknown",    False),
)

_STEP_IDS: tuple[str, ...] = tuple(s[0] for s in _STEPS)
_PAGE_TYPES: tuple[str, ...] = ("cover", "toc", "chapter", "content", "ending")
_PAGE_TYPE_FILES: dict[str, str] = {
    "cover": "01_cover.svg",
    "toc": "02_toc.svg",
    "chapter": "02_chapter.svg",
    "content": "03_content.svg",
    "ending": "04_ending.svg",
}
_PAGE_TYPE_LABELS: dict[str, str] = {
    "cover": "封面",
    "toc": "目录",
    "chapter": "章节",
    "content": "内容",
    "ending": "结束",
}
_INLINE_IMAGE_HREF_RE = re.compile(
    r'(?P<prefix>\b(?:href|xlink:href)=["\'])'
    r'data:(?P<mime>image/[^;\'"]+);base64,(?P<data>[^\'"]+)'
    r'(?P<suffix>["\'])',
    re.IGNORECASE,
)
_SVG_ROOT_RE = re.compile(r"<svg\b(?P<attrs>[^>]*)>(?P<body>[\s\S]*?)</svg>", re.IGNORECASE)
_SVG_TEXT_RE = re.compile(r"<text\b(?P<attrs>[^>]*)>(?P<body>[\s\S]*?)</text>", re.IGNORECASE)
_SVG_ATTR_RE = re.compile(r'([:\w-]+)\s*=\s*("([^"]*)"|\'([^\']*)\')')
_PLACEHOLDER_TOKEN_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")
_IMPORTER_VERSION = "template_import.v2"
_EDITABLE_TEMPLATE_DECK_VERSION = "editable-template-deck.v4"


# ─────────────────────────────────────────────────────────────────────────────
# Helpers — state record management
# ─────────────────────────────────────────────────────────────────────────────


def _initial_state(ctx: PipelineContext) -> ImportTaskState:
    """Build the initial ``ImportTaskState`` for a fresh import."""
    now = time.time()
    steps: list[StepRecord] = [
        {"id": sid, "label": label, "status": "pending"}
        for sid, label, _kind, _skip in _STEPS
    ]
    state: ImportTaskState = {
        "import_id": ctx.import_id,
        "status": "processing",
        "stage": _STEPS[0][0],
        "progress": 0.0,
        "message": "",
        "created_at": now,
        "updated_at": now,
        "importer_version": _IMPORTER_VERSION,
        "review_required": False,
        "template_id": None,
        "label": ctx.label,
        "source_file": str(ctx.pptx_path),
        "slide_count": 0,
        "export_mode": "",
        "error": None,
        "error_kind": None,
        "steps": steps,
        "theme_colors": [],
    }
    return state


def _ensure_state(ctx: PipelineContext) -> ImportTaskState:
    state = persistence.read_state(ctx.import_id)
    if state is None:
        state = _initial_state(ctx)
        persistence.write_state(ctx.import_id, state)
        return state
    # Make sure all six steps exist (safety net for partially migrated states).
    by_id = {s.get("id"): s for s in state.get("steps") or [] if isinstance(s, dict)}
    if set(_STEP_IDS) - set(by_id.keys()):
        steps: list[StepRecord] = []
        for sid, label, _kind, _skip in _STEPS:
            existing = by_id.get(sid)
            if existing:
                steps.append(existing)
            else:
                steps.append({"id": sid, "label": label, "status": "pending"})
        state["steps"] = steps
        persistence.write_state(ctx.import_id, state)
    return state


def _step_record(state: ImportTaskState, step_id: str) -> StepRecord:
    for record in state.get("steps") or []:
        if isinstance(record, dict) and record.get("id") == step_id:
            return record  # type: ignore[return-value]
    # Defensive fallback — append a fresh record.
    record = {"id": step_id, "label": step_id, "status": "pending"}
    steps = list(state.get("steps") or [])
    steps.append(record)
    state["steps"] = steps
    return record  # type: ignore[return-value]


def _set_step_status(
    state: ImportTaskState,
    step_id: str,
    status: str,
    *,
    started_at: float | None = None,
    ended_at: float | None = None,
    message: str | None = None,
    error: str | None = None,
) -> None:
    record = _step_record(state, step_id)
    record["status"] = status  # type: ignore[typeddict-item]
    if started_at is not None:
        record["started_at"] = started_at
    if ended_at is not None:
        record["ended_at"] = ended_at
        if "started_at" in record:
            try:
                duration = max(0.0, float(ended_at) - float(record["started_at"]))
                record["duration_ms"] = int(duration * 1000)
            except (TypeError, ValueError):
                pass
    if message is not None:
        record["message"] = message
    if error is not None:
        record["error"] = error


def _progress_for(step_id: str) -> float:
    if step_id not in _STEP_IDS:
        return 0.0
    idx = _STEP_IDS.index(step_id)
    # Map 0..len-1 → ~0.05..0.85; review step caps at 0.9 to leave room for confirm.
    span = max(1, len(_STEP_IDS))
    return min(0.9, 0.05 + (idx + 1) / span * 0.85)


def _result_from_state(state: ImportTaskState) -> ImportResult:
    return ImportResult(
        template_id=str(state.get("template_id") or ""),
        label=str(state.get("label") or ""),
        status=state.get("status") or "processing",  # type: ignore[arg-type]
        stage=str(state.get("stage") or _STEPS[0][0]),
        progress=float(state.get("progress") or 0.0),
        message=str(state.get("message") or ""),
        steps=list(state.get("steps") or []),
        review_required=bool(state.get("review_required")),
        export_mode=str(state.get("export_mode") or ""),
        slide_count=int(state.get("slide_count") or 0),
        theme_colors=list(state.get("theme_colors") or []),
        error=str(state.get("error") or ""),
        error_kind=state.get("error_kind"),
    )


# ─────────────────────────────────────────────────────────────────────────────
# Helpers — page-type / asset assembly
# ─────────────────────────────────────────────────────────────────────────────


def _load_or_init_review(ctx: PipelineContext, template_id: str) -> ReviewDraft:
    review = persistence.read_review(ctx.import_id)
    if review is not None:
        return review
    draft: ReviewDraft = {  # type: ignore[typeddict-unknown-key]
        "schema_version": 2,
        "import_id": ctx.import_id,
        "template_id": template_id,
        "label": ctx.label,
        "status": "processing",
        "slide_count": 0,
        "page_selections": {},
        "page_type_candidates": {},
        "assets": {},
        "assets_full": [],
        "preserve_texts": [],
        "placeholder_hints": {},
        "element_actions": [],
        "conversation": [],
        "feedback_history": [],
        "llm_trace": [],
    }
    return draft


def _candidate_to_full(candidate: AssetCandidate) -> dict[str, Any]:
    return {
        "asset_id": candidate.asset_id,
        "file_name": candidate.file_name,
        "pages": list(candidate.pages),
        "occurrences": [
            {
                "slide_index": occ.slide_index,
                "x": occ.x,
                "y": occ.y,
                "width": occ.width,
                "height": occ.height,
                "layer": occ.layer,
            }
            for occ in candidate.occurrences
        ],
        "sha1": candidate.sha1,
        "phash": candidate.phash,
        "bytes": candidate.bytes,
        "mime": candidate.mime,
        "width_px": candidate.width_px,
        "height_px": candidate.height_px,
        "position_stable": candidate.position_stable,
        "recommended_role": candidate.recommended_role,
        "role_source": candidate.role_source,
        "role_confidence": candidate.role_confidence,
    }


def _assets_role_overrides(candidates: list[AssetCandidate]) -> dict[str, Any]:
    """Build the trimmed ``review.assets`` dict of role overrides."""
    out: dict[str, Any] = {}
    for c in candidates:
        out[c.asset_id] = {
            "asset_id": c.asset_id,
            "name": c.file_name,
            "role": c.recommended_role,
            "role_source": c.role_source,
        }
    return out


def _slide_signal_view(
    metric: SlideMetrics,
    total: int,
    candidates: list[AssetCandidate],
) -> dict[str, Any]:
    """Build a minimal slide-record dict consumable by ``page_classifier.classify``."""
    images: list[dict[str, Any]] = []
    for c in candidates:
        for occ in c.occurrences:
            if occ.slide_index == metric.slide_index:
                images.append({"width": float(occ.width), "height": float(occ.height)})
    return {
        "slide_index": metric.slide_index,
        "total_slides": total,
        "text_runs": [],
        "text_density": 0.0,
        "images": images,
        "font_sizes": [],
    }


def _build_initial_manifest(
    ctx: PipelineContext,
    render_result: RenderResult | None,
) -> dict[str, Any]:
    """Best-effort initial manifest from ``manifest.build_manifest`` if available."""
    try:
        from .manifest import build_manifest

        out_dir = ctx.work_dir / "analysis"
        out_dir.mkdir(parents=True, exist_ok=True)
        return build_manifest(ctx.pptx_path, out_dir)
    except Exception:  # noqa: BLE001 — best-effort initial scan
        canvas = (1280, 720)
        if render_result is not None:
            canvas = render_result.canvas
        return {
            "source": {"name": Path(ctx.pptx_path).name},
            "slideSize": {"width_px": canvas[0], "height_px": canvas[1]},
            "theme": {"colors": {}, "fonts": {}},
            "assets": {"commonAssets": [], "allAssets": []},
            "slides": [],
        }


def _pptx_basic_info(pptx_path: Path) -> tuple[int, tuple[int, int]]:
    """Read slide count/canvas without requiring a render backend."""
    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        slide_count = len(prs.slides)
        # EMU → px at 96 DPI. 13.333in × 7.5in becomes 1280 × 720.
        width = int(round(int(prs.slide_width) / 9525)) or 1280
        height = int(round(int(prs.slide_height) / 9525)) or 720
        return slide_count, (width, height)
    except Exception:  # noqa: BLE001 - fallback to raw package inspection
        pass

    try:
        import zipfile

        with zipfile.ZipFile(pptx_path, "r") as zf:
            slide_count = sum(
                1
                for name in zf.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml$", name)
            )
        return slide_count, (1280, 720)
    except Exception:  # noqa: BLE001
        return 0, (1280, 720)


def _default_page_selections(slide_count: int) -> dict[str, int | None]:
    if slide_count <= 0:
        return {pt: None for pt in _PAGE_TYPES}
    selections: dict[str, int | None] = {
        "cover": 1,
        "toc": None,
        "chapter": None,
        "content": 1 if slide_count == 1 else 2,
        "ending": slide_count if slide_count > 1 else None,
    }
    if slide_count >= 3:
        selections["toc"] = 2
    if slide_count >= 4:
        selections["chapter"] = 3
        selections["content"] = 4
    return selections


def _pptist_deck_path(ctx: PipelineContext) -> Path:
    return ctx.work_dir / "pptist" / "deck.json"


def _pptist_source_deck_path(ctx: PipelineContext) -> Path:
    return ctx.work_dir / "pptist" / "source_deck.json"


def _load_pptist_deck(ctx: PipelineContext) -> dict[str, Any] | None:
    path = _pptist_deck_path(ctx)
    if not path.exists():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return None
    return data if isinstance(data, dict) else None


def _load_pptist_source_deck(ctx: PipelineContext) -> dict[str, Any] | None:
    path = _pptist_source_deck_path(ctx)
    if not path.exists():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return None
    return data if isinstance(data, dict) else None


def _agent_template_version(ctx: PipelineContext) -> str:
    agent_dir = ctx.work_dir / "agent_template"
    paths = [
        agent_dir / "manifest.patch.json",
        agent_dir / "design_spec.md",
        *[agent_dir / file_name for file_name in _PAGE_TYPE_FILES.values()],
    ]
    version = 0
    for path in paths:
        try:
            stat = path.stat()
        except OSError:
            continue
        version = max(version, stat.st_mtime_ns)
    return str(version)


def _agent_template_assets(ctx: PipelineContext) -> dict[str, bytes]:
    assets_dir = ctx.work_dir / "agent_template" / "assets"
    if not assets_dir.is_dir():
        return {}
    assets: dict[str, bytes] = {}
    for item in assets_dir.rglob("*"):
        if not item.is_file():
            continue
        try:
            data = item.read_bytes()
        except OSError:
            continue
        rel = item.relative_to(assets_dir).as_posix()
        assets[rel] = data
        assets[item.name] = data
    return assets


def _load_agent_template_svgs(ctx: PipelineContext) -> dict[str, str]:
    agent_dir = ctx.work_dir / "agent_template"
    if not agent_dir.is_dir():
        return {}
    if not any(
        (agent_dir / marker).exists()
        for marker in ("manifest.patch.json", "design_spec.md", "summary.md")
    ):
        return {}
    asset_map = _agent_template_assets(ctx)
    svgs: dict[str, str] = {}
    for page_type, file_name in _PAGE_TYPE_FILES.items():
        path = agent_dir / file_name
        if not path.exists():
            continue
        try:
            svg_text = path.read_text(encoding="utf-8")
        except OSError:
            continue
        if "<svg" not in svg_text or "</svg>" not in svg_text:
            continue
        inline = inline_svg_asset_refs(svg_text, asset_map)
        svgs[page_type] = _restore_pptist_brand_images(ctx, page_type, inline)
    return svgs


def _restore_pptist_brand_images(ctx: PipelineContext, page_type: str, svg_text: str) -> str:
    """Reinsert stable brand images if an Agent SVG accidentally dropped them."""

    deck = _load_pptist_deck(ctx)
    if not deck:
        return svg_text
    slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
    if not slides:
        return svg_text

    try:
        review = persistence.read_review(ctx.import_id)
    except Exception:  # noqa: BLE001 - preview restoration is best-effort
        review = {}
    page_selections = dict(review.get("page_selections") or {}) if isinstance(review, dict) else {}
    if not any(page_selections.get(pt) for pt in _PAGE_TYPES):
        page_selections = _default_page_selections(len(slides))
    slide_index = page_selections.get(page_type)
    if not slide_index:
        return svg_text
    try:
        slide = slides[int(slide_index) - 1]
    except (IndexError, TypeError, ValueError):
        return svg_text
    if not isinstance(slide, dict):
        return svg_text

    canvas_w = _as_float(deck.get("width"), 1280)
    canvas_h = _as_float(deck.get("height"), 720)
    restored: list[str] = []
    for raw in slide.get("elements") or []:
        if not isinstance(raw, dict) or raw.get("type") != "image" or not raw.get("src"):
            continue
        box = _element_bbox(raw)
        ratio = _area_ratio(box, canvas_w, canvas_h)
        is_brand_candidate = (
            "logo" in str(raw.get("name") or "").lower()
            or ratio <= 0.025
            or (ratio <= 0.12 and (box["y"] <= canvas_h * 0.24 or box["y"] + box["height"] >= canvas_h * 0.76))
            or (ratio <= 0.08 and _edge_element(box, canvas_w, canvas_h))
        )
        if not is_brand_candidate:
            continue
        src = str(raw.get("src") or "")
        element_id = str(raw.get("id") or "")
        if (src and src in svg_text) or (element_id and element_id in svg_text):
            continue
        markup = _svg_decorative_element(raw, canvas_w, canvas_h)
        if markup:
            restored.append(f'<g data-restored-brand="true">{markup}</g>')
    if not restored:
        return svg_text
    return svg_text.replace("</svg>", "\n  " + "\n  ".join(restored) + "\n</svg>", 1)


def _load_agent_design_spec(ctx: PipelineContext) -> str:
    for path in (
        ctx.work_dir / "agent_template" / "design_spec.md",
        ctx.work_dir / "agent_template" / "summary.md",
    ):
        if not path.exists():
            continue
        try:
            text = path.read_text(encoding="utf-8").strip()
        except OSError:
            continue
        if text:
            return text
    return ""


def _as_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _strip_html_text(value: Any) -> str:
    text = re.sub(r"<[^>]+>", " ", str(value or ""))
    return html.unescape(" ".join(text.split()))


def _first_inline_style(value: Any) -> dict[str, str]:
    raw = str(value or "")
    out: dict[str, str] = {}
    for match in re.finditer(r'style\s*=\s*["\']([^"\']+)["\']', raw, re.IGNORECASE):
        for item in match.group(1).split(";"):
            if ":" not in item:
                continue
            key, val = item.split(":", 1)
            out[key.strip().lower()] = val.strip()
    return out


def _pptist_text_style(el: dict[str, Any], box: dict[str, float]) -> dict[str, Any]:
    content = str(el.get("content") or "")
    style = _first_inline_style(content)
    color = style.get("color") or el.get("defaultColor") or "#111827"
    font = style.get("font-family") or el.get("defaultFontName") or "Arial"
    font = str(font).strip().strip("'\"") or "Arial"
    size_raw = str(style.get("font-size") or "")
    size_match = re.search(r"[-+]?\d+(?:\.\d+)?", size_raw)
    font_size = _as_float(size_match.group(0), 0.0) if size_match else 0.0
    if font_size <= 0:
        font_size = max(10.0, min(54.0, box["height"] * 0.42))
    align = style.get("text-align") or ""
    if not align:
        p_style = re.search(r"<p\b[^>]*style\s*=\s*[\"']([^\"']+)[\"']", content, re.IGNORECASE)
        if p_style:
            align_match = re.search(r"text-align\s*:\s*([^;]+)", p_style.group(1), re.IGNORECASE)
            if align_match:
                align = align_match.group(1).strip()
    align = align.lower()
    anchor = "middle" if align == "center" else "end" if align == "right" else "start"
    weight = str(style.get("font-weight") or "").lower()
    italic = str(style.get("font-style") or "").lower()
    return {
        "color": color,
        "font": font,
        "font_size": font_size,
        "anchor": anchor,
        "bold": weight in {"bold", "600", "700", "800", "900"},
        "italic": italic == "italic",
    }


def _pptist_rich_text_lines(
    value: Any,
    box: dict[str, float],
    *,
    default_color: str = "#111827",
    default_font: str = "Arial",
    default_weight: str = "400",
    align_hint: str = "",
) -> list[dict[str, Any]]:
    raw = str(value or "")
    blocks = re.findall(r"<p\b[^>]*>[\s\S]*?</p>", raw, flags=re.IGNORECASE)
    if not blocks:
        blocks = [raw]
    lines: list[dict[str, Any]] = []
    for block in blocks:
        text = _strip_html_text(block)
        if not text:
            continue
        style = _first_inline_style(block)
        size_raw = str(style.get("font-size") or "")
        size_match = re.search(r"[-+]?\d+(?:\.\d+)?", size_raw)
        font_size = _as_float(size_match.group(0), 0.0) if size_match else 0.0
        if font_size <= 0:
            font_size = max(10.0, min(54.0, box["height"] * 0.42))
        align = str(style.get("text-align") or align_hint or "").lower()
        anchor = "middle" if align in {"center", "middle"} else "end" if align == "right" else "start"
        weight = str(style.get("font-weight") or default_weight or "400").lower()
        font = str(style.get("font-family") or default_font or "Arial").strip().strip("'\"") or "Arial"
        lines.append(
            {
                "text": text,
                "color": style.get("color") or default_color or "#111827",
                "font": font,
                "font_size": font_size,
                "anchor": anchor,
                "weight": "700" if weight in {"bold", "600", "700", "800", "900"} else weight,
                "italic": str(style.get("font-style") or "").lower() == "italic",
            }
        )
    return lines


def _svg_attr(value: Any) -> str:
    return html.escape(str(value or ""), quote=True)


def _element_bbox(el: dict[str, Any]) -> dict[str, float]:
    return {
        "x": _as_float(el.get("left")),
        "y": _as_float(el.get("top")),
        "width": max(1.0, _as_float(el.get("width"), 1.0)),
        "height": max(1.0, _as_float(el.get("height"), 1.0)),
    }


def _pptist_transform_attr(el: dict[str, Any], box: dict[str, float]) -> str:
    rotate = _as_float(el.get("rotate"))
    flip_h = bool(el.get("flipH"))
    flip_v = bool(el.get("flipV"))
    if not rotate and not flip_h and not flip_v:
        return ""
    cx = box["x"] + box["width"] / 2
    cy = box["y"] + box["height"] / 2
    transforms: list[str] = []
    if rotate:
        transforms.append(f"rotate({rotate:.2f} {cx:.2f} {cy:.2f})")
    if flip_h or flip_v:
        sx = -1 if flip_h else 1
        sy = -1 if flip_v else 1
        transforms.append(
            f"translate({cx:.2f} {cy:.2f}) scale({sx} {sy}) translate({-cx:.2f} {-cy:.2f})"
        )
    return f' transform="{" ".join(transforms)}"'


def _edge_element(box: dict[str, float], canvas_w: float, canvas_h: float) -> bool:
    return (
        box["x"] < canvas_w * 0.14
        or box["y"] < canvas_h * 0.14
        or box["x"] + box["width"] > canvas_w * 0.86
        or box["y"] + box["height"] > canvas_h * 0.86
    )


def _area_ratio(box: dict[str, float], canvas_w: float, canvas_h: float) -> float:
    return box["width"] * box["height"] / max(1.0, canvas_w * canvas_h)


def _is_decorative_pptist_element(el: dict[str, Any], canvas_w: float, canvas_h: float) -> bool:
    box = _element_bbox(el)
    ratio = _area_ratio(box, canvas_w, canvas_h)
    name = str(el.get("name") or "").lower()
    el_type = el.get("type")
    if el_type == "image":
        image_type = str(el.get("imageType") or "")
        return (
            image_type == "background"
            or "logo" in name
            or ratio >= 0.72
            or ratio <= 0.025
            or (ratio <= 0.12 and (box["y"] <= canvas_h * 0.24 or box["y"] + box["height"] >= canvas_h * 0.76))
            or (ratio <= 0.08 and _edge_element(box, canvas_w, canvas_h))
        )
    if el_type == "text":
        content = _strip_html_text(el.get("content"))
        text_type = str(el.get("textType") or "")
        return (
            "{{" in content
            or text_type in {"header", "footer", "partNumber", "itemNumber"}
            or (ratio <= 0.035 and _edge_element(box, canvas_w, canvas_h))
        )
    if el_type == "shape":
        text = el.get("text")
        text_content = _strip_html_text(text.get("content")) if isinstance(text, dict) else ""
        if "{{" in text_content:
            return True
        if text_content:
            return False
        return ratio <= 0.5 or _edge_element(box, canvas_w, canvas_h)
    if el_type == "line":
        return True
    return False


def _text_element_box(el: dict[str, Any]) -> dict[str, float] | None:
    if el.get("type") == "text" and _strip_html_text(el.get("content")):
        return _element_bbox(el)
    if el.get("type") == "shape" and isinstance(el.get("text"), dict):
        if _strip_html_text(el["text"].get("content")):
            return _element_bbox(el)
    return None


def _pick_box(
    boxes: list[dict[str, float]],
    canvas_w: float,
    canvas_h: float,
    *,
    default: tuple[float, float, float, float],
    top_bias: bool = False,
) -> dict[str, float]:
    candidates = boxes
    if top_bias:
        top_candidates = [box for box in boxes if box["y"] <= canvas_h * 0.36]
        if top_candidates:
            candidates = top_candidates
    if not candidates:
        x, y, w, h = default
        return {"x": x, "y": y, "width": w, "height": h}
    return max(candidates, key=lambda box: box["width"] * box["height"])


def _placeholder_text(
    token: str,
    box: dict[str, float],
    *,
    font_size: float = 28,
    weight: str = "700",
    color: str = "#111827",
    anchor: str = "start",
) -> str:
    x = box["x"] + (box["width"] / 2 if anchor == "middle" else 0)
    y = box["y"] + box["height"] * 0.62
    return (
        f'<text x="{x:.2f}" y="{y:.2f}" text-anchor="{anchor}" '
        f'fill="{_svg_attr(color)}" font-family="Arial, sans-serif" '
        f'font-size="{font_size:.0f}" font-weight="{_svg_attr(weight)}">'
        f'{{{{{_svg_attr(token)}}}}}</text>'
    )


def _content_area_markup(box: dict[str, float]) -> str:
    cx = box["x"] + box["width"] / 2
    cy = box["y"] + box["height"] / 2
    return f"""
  <g id="content-area">
    <rect x="{box['x']:.2f}" y="{box['y']:.2f}" width="{box['width']:.2f}" height="{box['height']:.2f}" fill="none" stroke="#94A3B8" stroke-width="2" stroke-dasharray="8 8" opacity="0.75"/>
    <text x="{cx:.2f}" y="{cy:.2f}" text-anchor="middle" fill="#94A3B8" font-family="Arial, sans-serif" font-size="18">{{{{CONTENT_AREA}}}}</text>
  </g>"""


def _default_content_box(canvas_w: float, canvas_h: float) -> dict[str, float]:
    return {
        "x": canvas_w * 0.08,
        "y": canvas_h * 0.24,
        "width": canvas_w * 0.84,
        "height": canvas_h * 0.62,
    }


def _content_box_from_slide(slide: dict[str, Any], canvas_w: float, canvas_h: float) -> dict[str, float]:
    candidates: list[dict[str, float]] = []
    for raw in slide.get("elements") or []:
        if not isinstance(raw, dict):
            continue
        if _is_decorative_pptist_element(raw, canvas_w, canvas_h):
            continue
        box = _element_bbox(raw)
        if box["y"] < canvas_h * 0.16:
            continue
        candidates.append(box)
    if not candidates:
        return _default_content_box(canvas_w, canvas_h)
    return max(candidates, key=lambda box: box["width"] * box["height"])


def _svg_background(slide: dict[str, Any], canvas_w: float, canvas_h: float) -> str:
    background = slide.get("background")
    if not isinstance(background, dict):
        return f'<rect width="{canvas_w:.0f}" height="{canvas_h:.0f}" fill="#ffffff"/>'
    if background.get("type") == "image" and isinstance(background.get("image"), dict):
        src = background["image"].get("src")
        if src:
            return (
                f'<rect width="{canvas_w:.0f}" height="{canvas_h:.0f}" fill="#ffffff"/>'
                f'<image href="{_svg_attr(src)}" x="0" y="0" width="{canvas_w:.0f}" height="{canvas_h:.0f}" preserveAspectRatio="xMidYMid slice"/>'
            )
    color = background.get("color") or "#ffffff"
    return f'<rect width="{canvas_w:.0f}" height="{canvas_h:.0f}" fill="{_svg_attr(color)}"/>'


def _svg_outline(outline: Any) -> str:
    if not isinstance(outline, dict) or outline.get("style") == "none":
        return 'stroke="none"'
    color = outline.get("color") or "#111827"
    width = _as_float(outline.get("width"), 1.0)
    return f'stroke="{_svg_attr(color)}" stroke-width="{width:.2f}"'


def _pptist_shape_view_box(el: dict[str, Any]) -> tuple[float, float]:
    view_box = el.get("viewBox")
    if isinstance(view_box, list) and len(view_box) >= 2:
        return (
            max(1.0, _as_float(view_box[0], 1.0)),
            max(1.0, _as_float(view_box[1], 1.0)),
        )
    numbers = [
        _as_float(match.group(0), 0.0)
        for match in re.finditer(r"[-+]?\d+(?:\.\d+)?", str(el.get("path") or ""))
    ]
    xs = numbers[0::2]
    ys = numbers[1::2]
    if xs and ys:
        return max(1.0, max(xs) - min(xs)), max(1.0, max(ys) - min(ys))
    return 200.0, 200.0


def _svg_decorative_element(el: dict[str, Any], canvas_w: float, canvas_h: float) -> str:
    box = _element_bbox(el)
    transform = _pptist_transform_attr(el, box)
    el_type = el.get("type")
    if el_type == "image" and el.get("src"):
        return (
            f'<image href="{_svg_attr(el.get("src"))}" x="{box["x"]:.2f}" y="{box["y"]:.2f}" '
            f'width="{box["width"]:.2f}" height="{box["height"]:.2f}" preserveAspectRatio="none"{transform}/>'
        )
    if el_type == "shape" and el.get("path"):
        vb_w, vb_h = _pptist_shape_view_box(el)
        fill = el.get("fill") or "none"
        opacity = _as_float(el.get("opacity"), 1.0)
        return (
            f'<g{transform}>'
            f'<g transform="translate({box["x"]:.2f} {box["y"]:.2f}) scale({box["width"] / vb_w:.6f} {box["height"] / vb_h:.6f})">'
            f'<path d="{_svg_attr(el.get("path"))}" fill="{_svg_attr(fill)}" {_svg_outline(el.get("outline"))} opacity="{opacity:.3f}"/>'
            "</g>"
            "</g>"
        )
    if el_type == "line":
        start = el.get("start") if isinstance(el.get("start"), list) else [box["x"], box["y"]]
        end = el.get("end") if isinstance(el.get("end"), list) else [box["x"] + box["width"], box["y"]]
        color = el.get("color") or "#111827"
        return (
            f'<line x1="{_as_float(start[0]):.2f}" y1="{_as_float(start[1]):.2f}" '
            f'x2="{_as_float(end[0]):.2f}" y2="{_as_float(end[1]):.2f}" '
            f'stroke="{_svg_attr(color)}" stroke-width="2"{transform}/>'
        )
    if el_type == "text":
        content = _strip_html_text(el.get("content"))
        if not content:
            return ""
        text_style = _pptist_text_style(el, box)
        anchor = str(text_style["anchor"])
        x = box["x"] + (box["width"] / 2 if anchor == "middle" else box["width"] if anchor == "end" else 0)
        font_size = float(text_style["font_size"])
        weight = ' font-weight="700"' if text_style["bold"] else ""
        italic = ' font-style="italic"' if text_style["italic"] else ""
        return (
            f'<text x="{x:.2f}" y="{(box["y"] + font_size * 0.86):.2f}" text-anchor="{anchor}" '
            f'fill="{_svg_attr(text_style["color"])}" font-family="{_svg_attr(text_style["font"])}, sans-serif" '
            f'font-size="{font_size:.1f}"{weight}{italic}{transform}>'
            f'{html.escape(content)}</text>'
        )
    return ""


def _pptist_placeholders_for_page(
    slide: dict[str, Any],
    page_type: str,
    canvas_w: float,
    canvas_h: float,
) -> tuple[str, dict[str, float] | None]:
    text_boxes = [
        box
        for raw in slide.get("elements") or []
        if isinstance(raw, dict)
        for box in [_text_element_box(raw)]
        if box is not None and not _is_decorative_pptist_element(raw, canvas_w, canvas_h)
    ]
    title_box = _pick_box(
        text_boxes,
        canvas_w,
        canvas_h,
        default=(canvas_w * 0.12, canvas_h * 0.16, canvas_w * 0.76, canvas_h * 0.12),
        top_bias=True,
    )
    parts: list[str] = ['<g id="template-placeholders">']
    content_box: dict[str, float] | None = None
    if page_type == "cover":
        parts.append(_placeholder_text("TITLE", title_box, font_size=42, anchor="middle"))
        subtitle = {
            "x": canvas_w * 0.22,
            "y": min(canvas_h * 0.58, title_box["y"] + title_box["height"] + canvas_h * 0.08),
            "width": canvas_w * 0.56,
            "height": canvas_h * 0.08,
        }
        parts.append(_placeholder_text("SUBTITLE", subtitle, font_size=22, weight="400", color="#4B5563", anchor="middle"))
        meta_y = min(canvas_h * 0.82, subtitle["y"] + subtitle["height"] + canvas_h * 0.08)
        for idx, name in enumerate(("AUTHOR", "DATE")):
            box = {
                "x": canvas_w * 0.24,
                "y": meta_y + idx * canvas_h * 0.055,
                "width": canvas_w * 0.52,
                "height": canvas_h * 0.045,
            }
            parts.append(_placeholder_text(name, box, font_size=15, weight="600", color="#64748B", anchor="middle"))
    elif page_type == "toc":
        toc_box = _content_box_from_slide(slide, canvas_w, canvas_h)
        item_h = max(36.0, min(54.0, toc_box["height"] / 6.2))
        gap = max(10.0, item_h * 0.28)
        for idx in range(5):
            box = {
                "x": toc_box["x"],
                "y": toc_box["y"] + idx * (item_h + gap),
                "width": toc_box["width"],
                "height": item_h,
            }
            parts.append(_placeholder_text(f"TOC_ITEM_{idx + 1}", box, font_size=22, weight="600", color="#374151"))
    elif page_type == "chapter":
        num_box = {
            "x": canvas_w * 0.12,
            "y": canvas_h * 0.34,
            "width": canvas_w * 0.18,
            "height": canvas_h * 0.22,
        }
        chapter_box = {
            "x": canvas_w * 0.28,
            "y": canvas_h * 0.42,
            "width": canvas_w * 0.58,
            "height": canvas_h * 0.12,
        }
        parts.append(_placeholder_text("CHAPTER_NUMBER", num_box, font_size=72, color="#94A3B8", anchor="middle"))
        parts.append(_placeholder_text("CHAPTER_TITLE", chapter_box, font_size=38))
    elif page_type == "ending":
        ending_title = {
            "x": canvas_w * 0.16,
            "y": canvas_h * 0.36,
            "width": canvas_w * 0.68,
            "height": canvas_h * 0.12,
        }
        ending_message = {
            "x": canvas_w * 0.22,
            "y": canvas_h * 0.52,
            "width": canvas_w * 0.56,
            "height": canvas_h * 0.08,
        }
        parts.append(_placeholder_text("ENDING_TITLE", ending_title, font_size=42, anchor="middle"))
        parts.append(_placeholder_text("ENDING_MESSAGE", ending_message, font_size=22, weight="400", color="#64748B", anchor="middle"))
    else:
        parts.append(_placeholder_text("PAGE_TITLE", title_box, font_size=30))
        content_box = _content_box_from_slide(slide, canvas_w, canvas_h)
        parts.append(_content_area_markup(content_box))
    parts.append("</g>")
    return "\n  ".join(parts), content_box


def _pptist_slide_to_template_svg(
    slide: dict[str, Any],
    page_type: str,
    canvas_w: int,
    canvas_h: int,
) -> tuple[str, dict[str, float] | None]:
    decorative: list[str] = []
    for raw in slide.get("elements") or []:
        if not isinstance(raw, dict):
            continue
        if _is_decorative_pptist_element(raw, canvas_w, canvas_h):
            markup = _svg_decorative_element(raw, canvas_w, canvas_h)
            if markup:
                decorative.append(markup)
    placeholders, content_area = _pptist_placeholders_for_page(slide, page_type, canvas_w, canvas_h)
    body = "\n  ".join([_svg_background(slide, canvas_w, canvas_h), *decorative, placeholders])
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {canvas_w} {canvas_h}" width="{canvas_w}" height="{canvas_h}">\n'
        f"  {body}\n"
        "</svg>\n"
    ), content_area


def _shape_text_svg(el: dict[str, Any], box: dict[str, float], transform: str = "") -> str:
    raw_text = el.get("text")
    if not isinstance(raw_text, dict):
        return ""
    content = raw_text.get("content")
    lines = _pptist_rich_text_lines(
        content,
        box,
        default_color=str(raw_text.get("defaultColor") or "#111827"),
        default_font=str(raw_text.get("defaultFontName") or "Arial"),
        default_weight="600",
        align_hint=str(raw_text.get("align") or ""),
    )[:8]
    if not lines:
        return ""
    first = lines[0]
    anchor = str(first.get("anchor") or "start")
    x = box["x"] + (box["width"] / 2 if anchor == "middle" else box["width"] if anchor == "end" else box["width"] * 0.06)
    line_height = _as_float(raw_text.get("lineHeight"), 1.18)
    line_heights = [float(line.get("font_size") or 12.0) * max(1.0, line_height) for line in lines]
    total_height = sum(line_heights)
    y = box["y"] + max(0.0, (box["height"] - total_height) / 2.0) + float(first.get("font_size") or 12.0) * 0.86
    tspans: list[str] = []
    for index, line in enumerate(lines):
        dy = "0" if index == 0 else f"{line_heights[index - 1]:.2f}"
        attrs = [
            f'x="{x:.2f}"',
            f'dy="{dy}"',
            f'fill="{_svg_attr(line.get("color"))}"',
            f'font-family="{_svg_attr(line.get("font"))}, sans-serif"',
            f'font-size="{float(line.get("font_size") or 12.0):.1f}"',
            f'font-weight="{_svg_attr(line.get("weight") or "400")}"',
        ]
        if line.get("italic"):
            attrs.append('font-style="italic"')
        tspans.append(f'<tspan {" ".join(attrs)}>{html.escape(str(line.get("text") or "")[:240])}</tspan>')
    return (
        f'<text x="{x:.2f}" y="{y:.2f}" text-anchor="{anchor}" '
        f'fill="{_svg_attr(first.get("color"))}" font-family="{_svg_attr(first.get("font"))}, sans-serif" '
        f'font-size="{float(first.get("font_size") or 12.0):.1f}" font-weight="{_svg_attr(first.get("weight") or "400")}"{transform}>'
        f'{"".join(tspans)}</text>'
    )


def _svg_raw_pptist_element(el: dict[str, Any], canvas_w: float, canvas_h: float) -> str:
    box = _element_bbox(el)
    transform = _pptist_transform_attr(el, box)
    el_type = el.get("type")
    element_id = _svg_attr(el.get("id"))
    if el_type == "image" and el.get("src"):
        return (
            f'<image data-pptist-id="{element_id}" href="{_svg_attr(el.get("src"))}" '
            f'x="{box["x"]:.2f}" y="{box["y"]:.2f}" width="{box["width"]:.2f}" '
            f'height="{box["height"]:.2f}" preserveAspectRatio="none"{transform}/>'
        )
    if el_type == "shape" and el.get("path"):
        vb_w, vb_h = _pptist_shape_view_box(el)
        fill = el.get("fill") or "none"
        opacity = _as_float(el.get("opacity"), 1.0)
        shape = (
            f'<g data-pptist-id="{element_id}"{transform}>'
            f'<g transform="translate({box["x"]:.2f} {box["y"]:.2f}) scale({box["width"] / vb_w:.6f} {box["height"] / vb_h:.6f})">'
            f'<path d="{_svg_attr(el.get("path"))}" fill="{_svg_attr(fill)}" {_svg_outline(el.get("outline"))} opacity="{opacity:.3f}"/>'
            "</g>"
            f'{_shape_text_svg(el, box)}'
            "</g>"
        )
        return shape
    if el_type == "line":
        start = el.get("start") if isinstance(el.get("start"), list) else [box["x"], box["y"]]
        end = el.get("end") if isinstance(el.get("end"), list) else [box["x"] + box["width"], box["y"]]
        color = el.get("color") or "#111827"
        width = max(1.0, _as_float(el.get("width"), 2.0))
        return (
            f'<line data-pptist-id="{element_id}" x1="{_as_float(start[0]):.2f}" y1="{_as_float(start[1]):.2f}" '
            f'x2="{_as_float(end[0]):.2f}" y2="{_as_float(end[1]):.2f}" '
            f'stroke="{_svg_attr(color)}" stroke-width="{width:.2f}"{transform}/>'
        )
    if el_type == "text":
        content = _strip_html_text(el.get("content"))
        if not content:
            return ""
        text_style = _pptist_text_style(el, box)
        anchor = str(text_style["anchor"])
        x = box["x"] + (box["width"] / 2 if anchor == "middle" else box["width"] if anchor == "end" else 0)
        font_size = float(text_style["font_size"])
        y = box["y"] + font_size * 0.86
        weight = ' font-weight="700"' if text_style["bold"] else ""
        italic = ' font-style="italic"' if text_style["italic"] else ""
        return (
            f'<text data-pptist-id="{element_id}" x="{x:.2f}" y="{y:.2f}" text-anchor="{anchor}" '
            f'fill="{_svg_attr(text_style["color"])}" font-family="{_svg_attr(text_style["font"])}, sans-serif" '
            f'font-size="{font_size:.1f}"{weight}{italic}{transform}>{html.escape(content[:320])}</text>'
        )
    return ""


def _pptist_slide_to_raw_template_svg(slide: dict[str, Any], canvas_w: int, canvas_h: int) -> str:
    parts = [_svg_background(slide, canvas_w, canvas_h)]
    for raw in slide.get("elements") or []:
        if isinstance(raw, dict):
            markup = _svg_raw_pptist_element(raw, canvas_w, canvas_h)
            if markup:
                parts.append(markup)
    body = "\n  ".join(parts)
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {canvas_w} {canvas_h}" width="{canvas_w}" height="{canvas_h}">\n'
        f"  {body}\n"
        "</svg>\n"
    )


def _data_uri_extension(mime_type: str, payload: bytes) -> str:
    mime = mime_type.lower()
    if mime == "image/svg+xml" or payload.lstrip().startswith(b"<svg"):
        return "svg"
    if mime in {"image/png"} or payload.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if mime in {"image/jpeg", "image/jpg"} or payload.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if mime == "image/webp" or (payload.startswith(b"RIFF") and payload[8:12] == b"WEBP"):
        return "webp"
    if mime == "image/gif" or payload.startswith((b"GIF87a", b"GIF89a")):
        return "gif"
    return "bin"


def _normalize_template_pack_svgs(
    svgs: dict[str, str],
    assets: dict[str, bytes] | None = None,
) -> tuple[dict[str, str], dict[str, bytes], list[ManifestAssetEntry]]:
    """Rewrite inline image data URIs into ``assets/...`` refs before install."""

    normalized_assets: dict[str, bytes] = dict(assets or {})
    common_assets: list[ManifestAssetEntry] = []
    seen_assets: set[str] = set()

    def _replace(match: re.Match[str]) -> str:
        mime = match.group("mime")
        raw = html.unescape(match.group("data"))
        payload_text = "".join(raw.split())
        payload_text += "=" * (-len(payload_text) % 4)
        try:
            payload = base64.b64decode(payload_text, validate=False)
        except Exception:  # noqa: BLE001 - keep malformed data URI untouched
            return match.group(0)
        digest = hashlib.sha1(payload).hexdigest()
        ext = _data_uri_extension(mime, payload)
        filename = f"inline_{digest[:12]}.{ext}"
        normalized_assets.setdefault(filename, payload)
        if filename not in seen_assets:
            seen_assets.add(filename)
            common_assets.append(
                {
                    "asset_id": digest[:16],
                    "file_name": filename,
                    "role": "decoration",  # type: ignore[typeddict-item]
                    "pages": [],
                    "bbox_norm": {"x": 0.0, "y": 0.0, "width": 0.0, "height": 0.0},
                    "sha1": digest,
                    "bytes": len(payload),
                }
            )
        return f"{match.group('prefix')}assets/{filename}{match.group('suffix')}"

    normalized_svgs = {
        page_type: _INLINE_IMAGE_HREF_RE.sub(_replace, svg_text or "")
        for page_type, svg_text in svgs.items()
    }
    return normalized_svgs, normalized_assets, common_assets


def _svg_attrs(attrs: str) -> dict[str, str]:
    parsed: dict[str, str] = {}
    for match in _SVG_ATTR_RE.finditer(attrs or ""):
        parsed[match.group(1)] = match.group(3) if match.group(3) is not None else match.group(4) or ""
    return parsed


def _svg_plain_text(value: str) -> str:
    return html.unescape(" ".join(re.sub(r"<[^>]+>", " ", value or "").split()))


def _svg_placeholder_child_attrs(body: str, token: str) -> dict[str, str]:
    for match in re.finditer(
        r"<(?:tspan|span)\b(?P<attrs>[^>]*)>(?P<body>[\s\S]*?)</(?:tspan|span)>",
        body or "",
        re.IGNORECASE,
    ):
        if token in _svg_plain_text(match.group("body")):
            return _svg_attrs(match.group("attrs"))
    return {}


def _remove_placeholder_text_nodes(svg_text: str) -> str:
    return _SVG_TEXT_RE.sub(
        lambda match: "" if _PLACEHOLDER_TOKEN_RE.search(_svg_plain_text(match.group("body"))) else match.group(0),
        svg_text,
    )


def _svg_to_data_uri(svg_text: str) -> str:
    encoded = base64.b64encode(svg_text.encode("utf-8")).decode("ascii")
    return f"data:image/svg+xml;base64,{encoded}"


def _decode_svg_data_uri(src: str) -> str:
    if not src.startswith("data:image/svg+xml;base64,"):
        return ""
    payload = src.split(",", 1)[1]
    payload += "=" * (-len(payload) % 4)
    try:
        return base64.b64decode(payload, validate=False).decode("utf-8", errors="replace")
    except Exception:  # noqa: BLE001
        return ""


def _css_value(raw: str, name: str) -> str | None:
    match = re.search(rf"{re.escape(name)}\s*:\s*([^;\"']+)", raw, re.IGNORECASE)
    return match.group(1).strip() if match else None


def _svg_style_value(
    child_attrs: dict[str, str],
    parent_attrs: dict[str, str],
    attr_name: str,
    *css_names: str,
) -> str | None:
    for attrs in (child_attrs, parent_attrs):
        direct = attrs.get(attr_name)
        if direct:
            return direct.strip()
        style = attrs.get("style") or ""
        for css_name in css_names or (attr_name,):
            value = _css_value(style, css_name)
            if value:
                return value.strip()
    return None


def _source_text_style(el: dict[str, Any], canvas_w: int) -> dict[str, Any]:
    text_source: Any = el
    if el.get("type") == "shape" and isinstance(el.get("text"), dict):
        text_source = el["text"]
    raw_content = str(text_source.get("content") or "")
    color = (
        _css_value(raw_content, "color")
        or str(text_source.get("defaultColor") or el.get("defaultColor") or "#111827")
    )
    font = (
        _css_value(raw_content, "font-family")
        or str(text_source.get("defaultFontName") or el.get("defaultFontName") or "Arial")
    ).split(",")[0].strip().strip("'\"") or "Arial"
    size_raw = _css_value(raw_content, "font-size")
    font_size = _as_float(size_raw.replace("px", "") if size_raw else None, 0)
    if font_size <= 0:
        box = _element_bbox(el)
        font_size = max(10.0, min(72.0, box["height"] * 0.46))
    weight = _css_value(raw_content, "font-weight") or str(text_source.get("fontWeight") or "700")
    align = (
        _css_value(raw_content, "text-align")
        or str(text_source.get("align") or text_source.get("textAlign") or "")
    ).lower()
    if align not in {"center", "right", "left"}:
        box = _element_bbox(el)
        cx = box["x"] + box["width"] / 2
        align = "center" if abs(cx - canvas_w / 2) < canvas_w * 0.14 else "left"
    return {"color": color, "font": font, "font_size": font_size, "weight": weight, "align": align}


def _source_text_slots(
    source_slide: dict[str, Any] | None,
    canvas_w: int,
    canvas_h: int,
) -> list[dict[str, Any]]:
    if not isinstance(source_slide, dict):
        return []
    slots: list[dict[str, Any]] = []
    for raw in source_slide.get("elements") or []:
        if not isinstance(raw, dict):
            continue
        if _is_decorative_pptist_element(raw, canvas_w, canvas_h):
            continue
        box = _text_element_box(raw)
        if box is None:
            continue
        text = _strip_html_text(raw.get("content") or (raw.get("text") or {}).get("content"))
        if not text or "{{" in text:
            continue
        slots.append(
            {
                "box": box,
                "style": _source_text_style(raw, canvas_w),
                "text": text,
                "area": box["width"] * box["height"],
                "center_x": box["x"] + box["width"] / 2,
                "center_y": box["y"] + box["height"] / 2,
            }
        )
    return slots


def _slot_for_token(
    token: str,
    page_type: str,
    slots: list[dict[str, Any]],
    canvas_w: int,
    canvas_h: int,
) -> dict[str, Any] | None:
    if not slots:
        return None
    by_area = sorted(slots, key=lambda slot: slot["area"], reverse=True)
    by_y = sorted(slots, key=lambda slot: (slot["box"]["y"], slot["box"]["x"]))
    top_slots = [slot for slot in by_y if slot["box"]["y"] < canvas_h * 0.32]
    lower_slots = [slot for slot in by_y if slot["box"]["y"] >= canvas_h * 0.32]
    centered = sorted(slots, key=lambda slot: abs(slot["center_x"] - canvas_w / 2))

    if page_type == "cover":
        if token in {"{{TITLE}}", "{{SUBTITLE}}"}:
            pool = [slot for slot in by_area if slot["box"]["y"] < canvas_h * 0.64] or by_area
            return pool[0 if token == "{{TITLE}}" or len(pool) == 1 else min(1, len(pool) - 1)]
        if token in {"{{AUTHOR}}", "{{DATE}}"}:
            ordered = lower_slots or by_y
            offset = {"{{AUTHOR}}": 0, "{{DATE}}": 1}[token]
            return ordered[min(offset, len(ordered) - 1)]
    if page_type == "toc" and token.startswith("{{TOC_ITEM_"):
        ordered = [slot for slot in by_y if slot["box"]["x"] > canvas_w * 0.22] or by_y
        idx = max(0, min(4, int(re.search(r"\d+", token).group(0)) - 1 if re.search(r"\d+", token) else 0))
        return ordered[min(idx, len(ordered) - 1)]
    if page_type == "chapter":
        if token == "{{CHAPTER_TITLE}}":
            return by_area[0]
        if token == "{{CHAPTER_NUMBER}}":
            numeric = [slot for slot in by_area if re.search(r"\d+", str(slot.get("text") or ""))]
            return (numeric or centered)[0]
    if page_type == "content":
        if token == "{{PAGE_TITLE}}":
            return (top_slots or by_y)[0]
        return None
    if page_type == "ending":
        if token == "{{ENDING_TITLE}}":
            return by_area[0]
        if token == "{{ENDING_MESSAGE}}":
            return (lower_slots or by_area)[min(1, len(lower_slots or by_area) - 1)]
    return None


def _svg_text_placeholders_to_pptist_elements(
    svg_text: str,
    page_type: str,
    canvas_w: int,
    canvas_h: int,
    source_slide: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    elements: list[dict[str, Any]] = []
    source_slots = _source_text_slots(source_slide, canvas_w, canvas_h)
    for index, match in enumerate(_SVG_TEXT_RE.finditer(svg_text), start=1):
        body = match.group("body")
        plain = _svg_plain_text(body)
        token_match = _PLACEHOLDER_TOKEN_RE.search(plain)
        if not token_match:
            continue
        token = token_match.group(0)
        attrs = _svg_attrs(match.group("attrs"))
        child_attrs = _svg_placeholder_child_attrs(body, token)
        data_box = attrs.get("data-pptist-box") or ""
        numbers = [_as_float(part) for part in re.findall(r"-?\d+(?:\.\d+)?", data_box)]
        font_size_raw = _svg_style_value(child_attrs, attrs, "font-size")
        if font_size_raw:
            font_size_raw = font_size_raw.replace("px", "").strip()
        font_size = max(10.0, _as_float(font_size_raw, 28.0))
        anchor = attrs.get("text-anchor") or "start"
        has_explicit_box = False
        if len(numbers) >= 4:
            left, top, width, height = numbers[:4]
            has_explicit_box = width > 8 and height > 8
        else:
            x = _as_float(attrs.get("x"), canvas_w / 2)
            y = _as_float(attrs.get("y"), canvas_h / 2)
            width = max(font_size * max(4, len(token)) * 0.68, canvas_w * 0.16)
            height = max(font_size * 1.35, 28.0)
            if anchor == "middle":
                left = x - width / 2
            elif anchor == "end":
                left = x - width
            else:
                left = x
            top = y - height * 0.76
        left = max(0.0, min(float(left), canvas_w - 1.0))
        top = max(0.0, min(float(top), canvas_h - 1.0))
        width = max(24.0, min(float(width), canvas_w - left))
        height = max(20.0, min(float(height), canvas_h - top))
        min_width_font_size = min(font_size, 42.0)
        min_text_width = min(
            float(canvas_w),
            max(24.0, min_width_font_size * max(4, len(token)) * 0.76 + 12.0),
        )
        if width < min_text_width:
            old_left = left
            old_width = width
            if anchor == "middle":
                left = old_left + old_width / 2 - min_text_width / 2
            elif anchor == "end":
                left = old_left + old_width - min_text_width
            width = min_text_width
            left = max(0.0, min(left, canvas_w - width))
            width = max(24.0, min(width, canvas_w - left))
        min_text_height = max(20.0, font_size * 1.35)
        if height < min_text_height:
            height = min(min_text_height, canvas_h - top)
        color = _svg_style_value(child_attrs, attrs, "fill", "fill", "color")
        font = _svg_style_value(child_attrs, attrs, "font-family")
        weight = _svg_style_value(child_attrs, attrs, "font-weight")
        align = "center" if anchor == "middle" else "right" if anchor == "end" else "left"
        source_slot = _slot_for_token(token, page_type, source_slots, canvas_w, canvas_h)
        if source_slot is not None:
            style = source_slot["style"]
            color = color or str(style.get("color") or "")
            font = font or str(style.get("font") or "")
            if not font_size_raw and not has_explicit_box:
                font_size = max(10.0, _as_float(style.get("font_size"), font_size))
            weight = weight or str(style.get("weight") or "")
            if not attrs.get("text-anchor"):
                align = str(style.get("align") or align)
        color = color or "#111827"
        font = (font or "Arial").split(",")[0].strip().strip("'\"") or "Arial"
        weight = weight or "700"
        element_id = hashlib.sha1(f"{page_type}:{index}:{token}".encode("utf-8")).hexdigest()[:10]
        elements.append(
            {
                "id": f"tpl_{element_id}",
                "type": "text",
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "rotate": 0,
                "content": (
                    f'<p style="text-align: {align};">'
                    f'<span style="color: {_svg_attr(color)}; font-family: {_svg_attr(font)}; '
                    f'font-size: {font_size:.1f}px; font-weight: {_svg_attr(weight)}; white-space: nowrap;">'
                    f"{html.escape(token)}</span></p>"
                ),
                "defaultFontName": font,
                "defaultColor": color,
                "lineHeight": 1.12,
                "wordSpace": 0,
                "paragraphSpace": 0,
                "inset": [0, 0, 0, 0],
                "textType": "content",
            }
        )
    return elements


def _editable_template_deck_from_svgs(
    *,
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    svgs: dict[str, str],
    design_spec: str,
    canvas_w: int,
    canvas_h: int,
    source_deck: dict[str, Any] | None = None,
) -> dict[str, Any]:
    slides: list[dict[str, Any]] = []
    source_slides = source_deck.get("slides") if isinstance(source_deck, dict) and isinstance(source_deck.get("slides"), list) else []
    page_selections = review.get("page_selections") if isinstance(review.get("page_selections"), dict) else {}
    for page_type in _PAGE_TYPES:
        svg_text = svgs.get(page_type) or ""
        if not svg_text:
            continue
        source_index = page_selections.get(page_type)
        if not source_index:
            source_index = _PAGE_TYPES.index(page_type) + 1
        source_slide = None
        try:
            source_slide = source_slides[int(source_index) - 1]
        except (TypeError, ValueError, IndexError):
            source_slide = None
        background_svg = _remove_placeholder_text_nodes(svg_text)
        slides.append(
            {
                "id": f"template_{page_type}",
                "elements": _svg_text_placeholders_to_pptist_elements(
                    svg_text,
                    page_type,
                    canvas_w,
                    canvas_h,
                    source_slide if isinstance(source_slide, dict) else None,
                ),
                "background": {
                    "type": "image",
                    "image": {
                        "src": _svg_to_data_uri(background_svg),
                        "size": "contain",
                    },
                },
                "remark": _PAGE_TYPE_LABELS.get(page_type, page_type),
                "templatePageType": page_type,
            }
        )
    return {
        "title": str(review.get("label") or state.get("label") or ctx.import_id),
        "width": canvas_w,
        "height": canvas_h,
        "theme": None,
        "slides": slides,
        "source": {
            "kind": "templateImport",
            "id": ctx.import_id,
            "stage": "clean_template_edit",
            "template_page_types": list(_PAGE_TYPES),
            "agent_template_version": _agent_template_version(ctx),
            "editable_template_deck_version": _EDITABLE_TEMPLATE_DECK_VERSION,
        },
        "design_spec": design_spec,
        "updated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
    }


def _deck_is_clean_template_edit(deck: dict[str, Any] | None) -> bool:
    if not isinstance(deck, dict):
        return False
    source = deck.get("source")
    return isinstance(source, dict) and source.get("stage") == "clean_template_edit"


def _pptist_clean_template_slide_to_svg(slide: dict[str, Any], canvas_w: int, canvas_h: int) -> str:
    background = slide.get("background") if isinstance(slide.get("background"), dict) else {}
    image = background.get("image") if isinstance(background.get("image"), dict) else {}
    background_svg = _decode_svg_data_uri(str(image.get("src") or ""))
    element_markup = [
        markup
        for raw in slide.get("elements") or []
        if isinstance(raw, dict)
        for markup in [_svg_raw_pptist_element(raw, canvas_w, canvas_h)]
        if markup
    ]
    if background_svg and "</svg>" in background_svg:
        return background_svg.replace("</svg>", "\n  " + "\n  ".join(element_markup) + "\n</svg>", 1)
    return _pptist_slide_to_raw_template_svg(slide, canvas_w, canvas_h)


def build_editable_template_deck(import_id: str) -> dict[str, Any] | None:
    state = persistence.read_state(import_id)
    review = persistence.read_review(import_id)
    if state is None or review is None:
        return None
    ctx = _rebuild_context(import_id, state)
    if _mode_from_state(state) != "agent" or not _agent_template_complete(ctx):
        return None
    svgs = _load_agent_template_svgs(ctx)
    if not all(svgs.get(page_type) for page_type in _PAGE_TYPES):
        return None
    design_spec = _load_agent_design_spec(ctx)
    deck = _load_pptist_deck(ctx) or {}
    source_deck = _load_pptist_source_deck(ctx) or deck
    canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
    canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
    return _editable_template_deck_from_svgs(
        ctx=ctx,
        state=state,
        review=review,
        svgs=svgs,
        design_spec=design_spec,
        canvas_w=canvas_w,
        canvas_h=canvas_h,
        source_deck=source_deck,
    )


def _mode_from_state(state: ImportTaskState) -> str:
    mode = str(state.get("collaboration_mode") or "direct")
    return mode if mode in {"direct", "agent", "classic"} else "direct"


def _ordered_direct_page_selections(slide_count: int) -> dict[str, int]:
    if slide_count != 5:
        raise PersistenceError(
            "Direct import requires exactly 5 slides in order: cover, toc, chapter, content, ending.",
            step_id="confirm",
            error_kind="persistence",
        )
    return {"cover": 1, "toc": 2, "chapter": 3, "content": 4, "ending": 5}


def _all_work_assets(ctx: PipelineContext) -> dict[str, bytes]:
    assets: dict[str, bytes] = {}
    for base in (ctx.work_dir / "assets", ctx.work_dir / "work" / "assets"):
        if not base.is_dir():
            continue
        for item in base.rglob("*"):
            if not item.is_file():
                continue
            try:
                data = item.read_bytes()
            except OSError:
                continue
            rel = item.relative_to(base).as_posix()
            assets[rel] = data
            assets[item.name] = data
    return assets


def _direct_design_spec(review: ReviewDraft, manifest: TemplateManifest) -> str:
    explicit = str(review.get("design_spec_md") or "").strip()
    if explicit:
        return explicit
    base = _default_design_spec(review, manifest).rstrip()
    return (
        f"{base}\n\n"
        "## Direct Import Contract\n"
        "- This template was imported as a five-page ordered template pack.\n"
        "- Page order is fixed: cover, toc, chapter, content, ending.\n"
        "- No automatic templateization or placeholder rewriting was applied.\n"
    )


def _agent_template_complete(ctx: PipelineContext) -> bool:
    agent_dir = ctx.work_dir / "agent_template"
    if not (agent_dir / "manifest.patch.json").exists():
        return False
    return all((agent_dir / file_name).exists() for file_name in _PAGE_TYPE_FILES.values())


def _load_agent_manifest_patch(ctx: PipelineContext) -> dict[str, Any]:
    path = ctx.work_dir / "agent_template" / "manifest.patch.json"
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return data if isinstance(data, dict) else {}


def _manifest_warnings_from_agent_patch(patch: dict[str, Any]) -> list[ManifestWarning]:
    warnings: list[ManifestWarning] = []
    for idx, raw in enumerate(list(patch.get("warnings") or [])):
        if isinstance(raw, dict):
            warnings.append(
                {
                    "code": str(raw.get("code") or f"agent_warning_{idx + 1}"),
                    "message": str(raw.get("message") or raw.get("detail") or raw),
                }
            )
        elif raw:
            warnings.append({"code": f"agent_warning_{idx + 1}", "message": str(raw)})
    return warnings


def _copy_pptist_source_to_layout(ctx: PipelineContext, layout_dir: Path) -> None:
    """Persist the editable PPTist source next to the installed template pack."""

    source_dir = ctx.work_dir / "pptist"
    deck_path = source_dir / "deck.json"
    current_pptx = source_dir / "current.pptx"
    if not deck_path.exists() and not current_pptx.exists():
        return
    target_dir = layout_dir / "pptist"
    target_dir.mkdir(parents=True, exist_ok=True)
    for source, name in ((deck_path, "deck.json"), (current_pptx, "current.pptx")):
        if not source.exists():
            continue
        try:
            (target_dir / name).write_bytes(source.read_bytes())
        except OSError:
            logger.warning("Could not persist editable PPTist template source: %s", source)


def _install_template_pack(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    manifest: TemplateManifest,
    svgs: dict[str, str],
    design_spec: str,
    assets: dict[str, bytes] | None = None,
    import_trace: list[LLMTraceEntry] | None = None,
) -> ImportResult:
    svgs, normalized_assets, extracted_assets = _normalize_template_pack_svgs(svgs, assets)
    if extracted_assets:
        manifest = dict(manifest)  # type: ignore[assignment]
        existing_assets = list(manifest.get("common_assets") or [])
        existing_names = {
            str(item.get("file_name"))
            for item in existing_assets
            if isinstance(item, dict) and item.get("file_name")
        }
        manifest["common_assets"] = [
            *existing_assets,
            *[item for item in extracted_assets if item["file_name"] not in existing_names],
        ]  # type: ignore[typeddict-item]
    pack = LayoutPack(
        template_id=manifest["template_id"],
        label=manifest["label"],
        manifest=manifest,
        svgs=svgs,  # type: ignore[arg-type]
        design_spec=design_spec,
        assets=normalized_assets,
        import_trace=import_trace or [],
    )
    try:
        install_result = persistence.transactional_install_layout(manifest["template_id"], pack)
        persistence.update_user_index(manifest["template_id"], manifest["label"], manifest)
    except PersistenceError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise PersistenceError(
            str(exc) or exc.__class__.__name__,
            step_id="confirm",
            error_kind="persistence",
        ) from exc
    _copy_pptist_source_to_layout(ctx, install_result.layout_dir)

    state["status"] = "complete"
    state["stage"] = "complete"
    state["progress"] = 1.0
    state["message"] = "Template registered"
    state["review_required"] = False
    state["template_id"] = manifest["template_id"]
    state["label"] = manifest["label"]
    state["slide_count"] = int(manifest.get("slide_count") or 0)
    state["export_mode"] = str(manifest.get("export_mode") or "")
    state["error"] = None
    state["error_kind"] = None
    state["updated_at"] = time.time()
    review["status"] = "complete"  # type: ignore[typeddict-item]
    persistence.write_review(ctx.import_id, review)
    persistence.write_state(ctx.import_id, state)
    result = _result_from_state(state)
    if install_result.already_complete:
        result.message = "Template already installed."
    return result


def _direct_preview_from_pptist_deck(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    deck: dict[str, Any],
) -> dict[str, Any]:
    slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
    page_selections = _ordered_direct_page_selections(len(slides))
    canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
    canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "template_id": str(review.get("template_id") or state.get("template_id") or ctx.import_id),
        "label": str(review.get("label") or state.get("label") or deck.get("title") or ctx.import_id),
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": 5,
        "canvas": {"width": canvas_w, "height": canvas_h},
        "theme": {"colors": {}, "fonts": {}},
        "page_selections": page_selections,  # type: ignore[typeddict-item]
        "content_area": _default_content_box(canvas_w, canvas_h),
        "warnings": [],
        "export_mode": "direct-pptist",
        "importer_version": _IMPORTER_VERSION,
        "source_kind": "pptist_import",
        "import_mode": "direct",
        "clean_generation_ready": True,
        "direct_ordered_template": True,
    }
    out: dict[str, Any] = {
        "template_id": manifest["template_id"],
        "label": manifest["label"],
        "cover_svg": "",
        "toc_svg": "",
        "chapter_svg": "",
        "content_svg": "",
        "ending_svg": "",
        "design_spec": _direct_design_spec(review, manifest),
        "theme_colors": list(state.get("theme_colors") or []),
        "warnings": [
            {
                "code": "direct_ordered_template",
                "message": "Direct import uses the five uploaded slides as-is; no templateization was applied.",
            }
        ],
        "placeholders": {},
        "preview_source": "direct_ordered",
    }
    for page_type, slide_index in page_selections.items():
        slide = slides[slide_index - 1]
        out[f"{page_type}_svg"] = _pptist_slide_to_raw_template_svg(slide, canvas_w, canvas_h)
        out["placeholders"][page_type] = []
    return out


def _direct_preview_from_rendered(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    render_result: RenderResult,
) -> dict[str, Any]:
    page_selections = _ordered_direct_page_selections(len(render_result.svg_files))
    asset_map = _all_work_assets(ctx)
    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "template_id": str(review.get("template_id") or state.get("template_id") or ctx.import_id),
        "label": str(review.get("label") or state.get("label") or ctx.import_id),
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": 5,
        "canvas": {"width": int(render_result.canvas[0] or 1280), "height": int(render_result.canvas[1] or 720)},
        "theme": {"colors": {}, "fonts": {}},
        "page_selections": page_selections,  # type: ignore[typeddict-item]
        "content_area": _default_content_box(float(render_result.canvas[0] or 1280), float(render_result.canvas[1] or 720)),
        "warnings": [],
        "export_mode": str(state.get("export_mode") or render_result.export_mode or "direct-rendered"),
        "importer_version": _IMPORTER_VERSION,
        "source_kind": "pptist_import",
        "import_mode": "direct",
        "clean_generation_ready": True,
        "direct_ordered_template": True,
    }
    out: dict[str, Any] = {
        "template_id": manifest["template_id"],
        "label": manifest["label"],
        "cover_svg": "",
        "toc_svg": "",
        "chapter_svg": "",
        "content_svg": "",
        "ending_svg": "",
        "design_spec": _direct_design_spec(review, manifest),
        "theme_colors": list(state.get("theme_colors") or []),
        "warnings": [
            {
                "code": "direct_ordered_template",
                "message": "Direct import uses the five uploaded slides as-is; no templateization was applied.",
            }
        ],
        "placeholders": {},
        "preview_source": "direct_ordered",
    }
    for page_type, slide_index in page_selections.items():
        svg_path = ctx.work_dir / "svg" / f"slide-{slide_index:03d}.svg"
        if not svg_path.exists():
            svg_path = ctx.work_dir / "svg" / f"slide-{slide_index:02d}.svg"
        text = svg_path.read_text(encoding="utf-8")
        out[f"{page_type}_svg"] = inline_svg_asset_refs(text, asset_map)
        out["placeholders"][page_type] = []
    return out


def _confirm_direct_ordered_template(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    template_id: str,
) -> ImportResult:
    deck = _load_pptist_deck(ctx)
    assets: dict[str, bytes] = {}
    render_result = _reload_render_result(ctx)
    if render_result is not None:
        page_selections = _ordered_direct_page_selections(len(render_result.svg_files))
        canvas_w = int(render_result.canvas[0] or 1280)
        canvas_h = int(render_result.canvas[1] or 720)
        assets = _all_work_assets(ctx)
        svgs = {}
        for page_type, slide_index in page_selections.items():
            svg_path = ctx.work_dir / "svg" / f"slide-{slide_index:03d}.svg"
            if not svg_path.exists():
                svg_path = ctx.work_dir / "svg" / f"slide-{slide_index:02d}.svg"
            svgs[page_type] = inline_svg_asset_refs(svg_path.read_text(encoding="utf-8"), assets)
        export_mode = str(state.get("export_mode") or render_result.export_mode or "direct-rendered")
        label = str(review.get("label") or state.get("label") or template_id)
    elif deck is not None:
        slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
        page_selections = _ordered_direct_page_selections(len(slides))
        canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
        canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
        svgs = {
            page_type: _pptist_slide_to_raw_template_svg(slides[slide_index - 1], canvas_w, canvas_h)
            for page_type, slide_index in page_selections.items()
        }
        export_mode = "direct-pptist"
        label = str(review.get("label") or state.get("label") or deck.get("title") or template_id)
    else:
        raise PersistenceError(
            "Direct import requires exactly 5 saved PPTist slides, or rendered SVGs from a 5-slide PPTX.",
            step_id="confirm",
            error_kind="persistence",
        )

    review["page_selections"] = dict(page_selections)  # type: ignore[typeddict-item]
    draft = review.setdefault("draft", {})  # type: ignore[assignment]
    if isinstance(draft, dict):
        draft["page_selections"] = dict(page_selections)
    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "schema_version": 2,
        "template_id": template_id,
        "label": label,
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": 5,
        "canvas": {"width": canvas_w, "height": canvas_h},
        "theme": {"colors": {}, "fonts": {}},
        "page_type_candidates": {pt: [idx] for pt, idx in page_selections.items()},  # type: ignore[typeddict-item]
        "page_selections": page_selections,  # type: ignore[typeddict-item]
        "common_assets": [],
        "content_area": _default_content_box(canvas_w, canvas_h),
        "warnings": [
            {
                "code": "direct_ordered_template",
                "message": "Direct import installed the five uploaded slides as-is; no templateization was applied.",
            }
        ],
        "imported_at": time.time(),
        "export_mode": export_mode,  # type: ignore[typeddict-item]
        "importer_version": _IMPORTER_VERSION,
        "source_kind": "pptist_import",
        "import_mode": "direct",
        "clean_generation_ready": True,
        "direct_ordered_template": True,
    }
    if not str(review.get("design_spec_md") or "").strip():
        raise PersistenceError(
            "Direct import requires generated design_spec.md before confirming.",
            step_id="confirm",
            error_kind="persistence",
        )
    design_spec = _direct_design_spec(review, manifest)
    return _install_template_pack(ctx, state, review, manifest, svgs, design_spec, assets)


def _preview_agent_template(ctx: PipelineContext, state: ImportTaskState, review: ReviewDraft) -> dict[str, Any]:
    out: dict[str, Any] = {
        "template_id": str(review.get("template_id") or state.get("template_id") or ctx.import_id),
        "label": str(review.get("label") or state.get("label") or ctx.import_id),
        "cover_svg": "",
        "toc_svg": "",
        "chapter_svg": "",
        "content_svg": "",
        "ending_svg": "",
        "design_spec": "",
        "theme_colors": list(state.get("theme_colors") or []),
        "warnings": [],
        "placeholders": {},
        "preview_source": "agent_artifact" if _agent_template_complete(ctx) else "agent_pending",
    }
    if not _agent_template_complete(ctx):
        out["warnings"].append(
            {
                "code": "agent_template_not_ready",
                "message": "Agent templateization has not produced a complete five-page template pack yet.",
            }
        )
        return out
    agent_svgs = _load_agent_template_svgs(ctx)
    for page_type in _PAGE_TYPES:
        out[f"{page_type}_svg"] = agent_svgs.get(page_type, "")
        out["placeholders"][page_type] = []
    design_spec = _load_agent_design_spec(ctx)
    out["design_spec"] = design_spec
    out["warnings"].append(
        {
            "code": "agent_template_preview",
            "message": "Preview includes Agent-authored template SVG outputs.",
        }
    )
    return out


def _confirm_agent_template(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    template_id: str,
) -> ImportResult:
    if not _agent_template_complete(ctx):
        raise PersistenceError(
            "Agent templateization has not produced a complete five-page template pack yet.",
            step_id="confirm",
            error_kind="persistence",
        )
    patch = _load_agent_manifest_patch(ctx)
    deck = _load_pptist_deck(ctx) or {}
    canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
    canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
    design_spec = str(deck.get("design_spec") or "").strip() if _deck_is_clean_template_edit(deck) else ""
    if _deck_is_clean_template_edit(deck):
        slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
        if len(slides) < len(_PAGE_TYPES):
            raise PersistenceError(
                "Edited Agent template deck must contain cover, toc, chapter, content, and ending slides.",
                step_id="confirm",
                error_kind="persistence",
            )
        svgs = {
            page_type: _pptist_clean_template_slide_to_svg(slides[index], canvas_w, canvas_h)
            for index, page_type in enumerate(_PAGE_TYPES)
            if isinstance(slides[index], dict)
        }
        export_mode = "agent-template-pptist"
    else:
        svgs = _load_agent_template_svgs(ctx)
        export_mode = "agent-template"
    missing = [pt for pt in _PAGE_TYPES if not svgs.get(pt)]
    if missing:
        raise PersistenceError(
            f"Agent template pack is missing SVGs for: {', '.join(missing)}",
            step_id="confirm",
            error_kind="persistence",
        )
    if not design_spec:
        design_spec = _load_agent_design_spec(ctx)
    if not design_spec:
        raise PersistenceError(
            "Agent templateization must produce design_spec.md before confirming.",
            step_id="confirm",
            error_kind="persistence",
        )
    page_selections = dict(review.get("page_selections") or patch.get("page_selections") or {})
    if not any(page_selections.get(pt) for pt in _PAGE_TYPES):
        slide_count = int(state.get("slide_count") or 5)
        page_selections = _default_page_selections(slide_count)
    content_area = patch.get("content_area") if isinstance(patch.get("content_area"), dict) else None
    if not content_area:
        content_area = _default_content_box(canvas_w, canvas_h)
    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "schema_version": 2,
        "template_id": template_id,
        "label": str(review.get("label") or state.get("label") or template_id),
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": int(state.get("slide_count") or len(deck.get("slides") or []) or 5),
        "canvas": {"width": canvas_w, "height": canvas_h},
        "theme": patch.get("theme") if isinstance(patch.get("theme"), dict) else {"colors": {}, "fonts": {}},
        "page_type_candidates": {pt: ([int(v)] if v else []) for pt, v in page_selections.items() if pt in _PAGE_TYPES},  # type: ignore[typeddict-item]
        "page_selections": {pt: int(v) if v is not None else None for pt, v in page_selections.items() if pt in _PAGE_TYPES},  # type: ignore[typeddict-item]
        "common_assets": [],
        "content_area": {
            "x": float(content_area.get("x", 0.08 * canvas_w)),
            "y": float(content_area.get("y", 0.24 * canvas_h)),
            "width": float(content_area.get("width", 0.84 * canvas_w)),
            "height": float(content_area.get("height", 0.62 * canvas_h)),
        },
        "warnings": _manifest_warnings_from_agent_patch(patch),
        "imported_at": time.time(),
        "export_mode": export_mode,  # type: ignore[typeddict-item]
        "importer_version": _IMPORTER_VERSION,
        "source_kind": "pptist_import",
        "import_mode": "agent",
        "clean_generation_ready": True,
    }
    assets = _agent_template_assets(ctx)
    return _install_template_pack(
        ctx,
        state,
        review,
        manifest,
        svgs,
        design_spec,
        assets,
        list(review.get("llm_trace") or []),
    )


def _confirm_from_pptist_deck(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    template_id: str,
    deck: dict[str, Any],
) -> ImportResult:
    slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
    if not slides:
        raise PersistenceError(
            "PPTist deck has no slides; save the deck before confirming.",
            step_id="confirm",
            error_kind="persistence",
        )

    canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
    canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
    page_selections = dict(review.get("page_selections") or {})
    if not any(page_selections.get(pt) for pt in _PAGE_TYPES):
        page_selections = _default_page_selections(len(slides))

    svgs: dict[str, str] = {}
    content_area_box: dict[str, float] | None = None
    for page_type in _PAGE_TYPES:
        slide_index = page_selections.get(page_type)
        if not slide_index:
            continue
        try:
            slide = slides[int(slide_index) - 1]
        except (IndexError, TypeError, ValueError):
            continue
        if not isinstance(slide, dict):
            continue
        svg_text, content_area = _pptist_slide_to_template_svg(slide, page_type, canvas_w, canvas_h)
        svgs[page_type] = svg_text
        if page_type == "content" and content_area is not None:
            content_area_box = content_area

    if str(state.get("collaboration_mode") or "") == "agent":
        agent_svgs = _load_agent_template_svgs(ctx)
        if agent_svgs:
            svgs.update(agent_svgs)

    if not svgs:
        raise PersistenceError(
            "No selected PPTist slides can be templateized.",
            step_id="confirm",
            error_kind="persistence",
        )

    if content_area_box is None:
        content_area_box = _default_content_box(canvas_w, canvas_h)

    page_selections_clean = {
        pt: int(v) if v is not None else None
        for pt, v in page_selections.items()
        if pt in _PAGE_TYPES
    }
    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "schema_version": 2,
        "template_id": template_id,
        "label": str(review.get("label") or state.get("label") or deck.get("title") or template_id),
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": len(slides),
        "canvas": {"width": canvas_w, "height": canvas_h},
        "theme": {"colors": {}, "fonts": {}},
        "page_type_candidates": {
            pt: ([int(v)] if v else [])
            for pt, v in page_selections_clean.items()
        },  # type: ignore[typeddict-item]
        "page_selections": page_selections_clean,  # type: ignore[typeddict-item]
        "common_assets": [],
        "content_area": {
            "x": float(content_area_box["x"]),
            "y": float(content_area_box["y"]),
            "width": float(content_area_box["width"]),
            "height": float(content_area_box["height"]),
        },
        "warnings": [
            {
                "code": "pptist_native_template",
                "message": "Template was rebuilt from PPTist deck JSON because no server SVG renderer was available.",
            }
        ],
        "imported_at": time.time(),
        "export_mode": "pptist-native",  # type: ignore[typeddict-item]
        "importer_version": _IMPORTER_VERSION,
        "source_kind": "pptist_import",
        "import_mode": str(state.get("collaboration_mode") or "direct"),
        "clean_generation_ready": True,
    }

    design_spec = str(review.get("design_spec_md") or "").strip()
    if not design_spec and str(state.get("collaboration_mode") or "") == "agent":
        design_spec = _load_agent_design_spec(ctx)
    if not design_spec:
        design_spec = _default_design_spec(review, manifest)

    pack = LayoutPack(
        template_id=template_id,
        label=manifest["label"],
        manifest=manifest,
        svgs=svgs,  # type: ignore[arg-type]
        design_spec=design_spec,
        assets={},
        import_trace=list(review.get("llm_trace") or []),
    )
    persistence.transactional_install_layout(template_id, pack)
    persistence.update_user_index(template_id, manifest["label"], manifest)

    now = time.time()
    state["status"] = "complete"
    state["stage"] = "review"
    state["progress"] = 1.0
    state["message"] = "Template registered"
    state["review_required"] = False
    state["template_id"] = template_id
    state["label"] = manifest["label"]
    state["slide_count"] = len(slides)
    state["export_mode"] = "pptist-native"
    state["error"] = None
    state["error_kind"] = None
    state["updated_at"] = now
    review["status"] = "complete"  # type: ignore[typeddict-item]
    persistence.write_review(ctx.import_id, review)
    persistence.write_state(ctx.import_id, state)
    return _result_from_state(state)


def _preview_from_pptist_deck(
    ctx: PipelineContext,
    state: ImportTaskState,
    review: ReviewDraft,
    deck: dict[str, Any],
) -> dict[str, Any]:
    slides = deck.get("slides") if isinstance(deck.get("slides"), list) else []
    if not slides:
        raise PersistenceError(
            "PPTist deck has no slides; save the deck before previewing.",
            step_id="preview",
            error_kind="persistence",
        )

    canvas_w = int(round(_as_float(deck.get("width"), 1280))) or 1280
    canvas_h = int(round(_as_float(deck.get("height"), 720))) or 720
    draft = review.get("draft") if isinstance(review.get("draft"), dict) else {}
    page_selections = dict(review.get("page_selections") or draft.get("page_selections") or {})
    if not any(page_selections.get(pt) for pt in _PAGE_TYPES):
        page_selections = _default_page_selections(len(slides))

    manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
        "template_id": str(review.get("template_id") or state.get("template_id") or ctx.import_id),
        "label": str(review.get("label") or state.get("label") or deck.get("title") or ctx.import_id),
        "source_file": str(state.get("source_file") or ctx.pptx_path.name),
        "slide_count": len(slides),
        "export_mode": "pptist-native",  # type: ignore[typeddict-item]
        "page_selections": {
            pt: int(value) if value is not None else None
            for pt, value in page_selections.items()
            if pt in _PAGE_TYPES
        },  # type: ignore[typeddict-item]
    }
    out: dict[str, Any] = {
        "template_id": manifest["template_id"],
        "label": manifest["label"],
        "cover_svg": "",
        "toc_svg": "",
        "chapter_svg": "",
        "content_svg": "",
        "ending_svg": "",
        "design_spec": str(review.get("design_spec_md") or _default_design_spec(review, manifest)),
        "theme_colors": list(state.get("theme_colors") or []),
        "warnings": [
            {
                "code": "pptist_native_preview",
                "message": "Preview was rebuilt from the saved PPTist deck because server SVG rendering is unavailable.",
            }
        ],
        "placeholders": {},
    }

    for page_type in _PAGE_TYPES:
        slide_index = page_selections.get(page_type)
        if not slide_index:
            continue
        try:
            slide = slides[int(slide_index) - 1]
        except (IndexError, TypeError, ValueError):
            continue
        if not isinstance(slide, dict):
            continue
        svg_text, content_area = _pptist_slide_to_template_svg(slide, page_type, canvas_w, canvas_h)
        out[f"{page_type}_svg"] = svg_text
        if page_type == "content" and content_area:
            out["placeholders"][page_type] = [
                {
                    "name": "CONTENT_AREA",
                    "bbox": {
                        "x": float(content_area["x"]),
                        "y": float(content_area["y"]),
                        "width": float(content_area["width"]),
                        "height": float(content_area["height"]),
                    },
                    "style": {},
                }
            ]
        else:
            out["placeholders"][page_type] = []

    if str(state.get("collaboration_mode") or "") == "agent":
        agent_svgs = _load_agent_template_svgs(ctx)
        if agent_svgs:
            for page_type, svg_text in agent_svgs.items():
                out[f"{page_type}_svg"] = svg_text
            agent_spec = _load_agent_design_spec(ctx)
            if agent_spec:
                out["design_spec"] = agent_spec
            out["warnings"].append(
                {
                    "code": "agent_template_preview",
                    "message": "Preview includes Agent-authored template SVG outputs.",
                }
            )

    return out


# ─────────────────────────────────────────────────────────────────────────────
# Helpers — inline_svg_asset_refs (preview = final-output parity)
# ─────────────────────────────────────────────────────────────────────────────


_HREF_PATTERN = re.compile(
    r'(?P<attr>(?:xlink:)?href)=(?P<quote>["\'])(?P<href>[^"\']*)(?P=quote)'
)
_DATA_URI_MAX_BYTES = 900_000


def _data_uri_for_bytes(name: str, data: bytes) -> str:
    if not data or len(data) > _DATA_URI_MAX_BYTES:
        return ""
    mime = mimetypes.guess_type(name)[0] or "application/octet-stream"
    encoded = base64.b64encode(data).decode("ascii")
    return f"data:{mime};base64,{encoded}"


def inline_svg_asset_refs(svg_text: str, asset_map: dict[str, bytes]) -> str:
    """Inline ``<image href>`` / ``xlink:href`` references to ``data:`` URIs.

    * ``data:``, ``http://``, ``https://`` and ``#`` references pass
      through unchanged (Requirement 12.2 — never produce ``http(s)://``).
    * Local references whose basename matches a key in ``asset_map`` are
      replaced with a ``data:<mime>;base64,...`` URI built from the
      corresponding bytes. References whose basename is *not* in the map
      are left untouched so they can be resolved at render time against
      the on-disk asset directory.
    * Everything outside the matched ``href="..."`` / ``xlink:href="..."``
      attribute values is left byte-identical, so ``{{NAME}}`` placeholder
      positions (line + column) are preserved across repeated invocations
      (Requirement 3.3, 12.2; Property 22).
    * Idempotent: ``inline(inline(x)) == inline(x)`` because already-inlined
      ``data:`` references short-circuit the substitution.
    """
    if not svg_text:
        return svg_text

    def _repl(match: re.Match[str]) -> str:
        attr = match.group("attr")
        quote = match.group("quote")
        href = match.group("href")
        if not href:
            return match.group(0)
        if href.startswith(("data:", "http://", "https://", "#")):
            return match.group(0)
        # Skip placeholder tokens — let the runtime templater fill them in.
        if "{{" in href and "}}" in href:
            return match.group(0)
        clean = href.split("#", 1)[0].split("?", 1)[0]
        if not clean:
            return match.group(0)
        basename = PurePosixPath(clean).name or clean
        payload = asset_map.get(basename)
        if payload is None:
            payload = asset_map.get(clean)
        if payload is None:
            return match.group(0)
        data_uri = _data_uri_for_bytes(basename, payload)
        if not data_uri:
            return match.group(0)
        return f"{attr}={quote}{data_uri}{quote}"

    return _HREF_PATTERN.sub(_repl, svg_text)


# ─────────────────────────────────────────────────────────────────────────────
# Pipeline
# ─────────────────────────────────────────────────────────────────────────────


class Pipeline:
    """Runs the template-import state machine for a single ``import_id``."""

    def __init__(self, ctx: PipelineContext | None = None) -> None:
        self._ctx = ctx

    # ── run ────────────────────────────────────────────────────────────

    async def run(self, ctx: PipelineContext) -> ImportResult:
        """Execute all pipeline steps in order until ``review_required``."""
        self._ctx = ctx
        state = _ensure_state(ctx)
        # Honor an already-complete state — nothing to redo.
        if state.get("status") == "complete":
            return _result_from_state(state)

        # Reset error markers for this run; per-step reruns clear earlier
        # error/skipped statuses inside retry_step before calling run().
        state["error"] = None
        state["error_kind"] = None
        if state.get("status") == "error":
            state["status"] = "processing"
        persistence.write_state(ctx.import_id, state)

        # Per-run intermediates carried between steps.
        render_result: RenderResult | None = None
        manifest_data: dict[str, Any] = {}
        candidates: list[AssetCandidate] = []
        warnings_aggregate: list[ManifestWarning] = []
        chrome_items: list[ChromeTextItem] = []
        review = _load_or_init_review(ctx, state.get("template_id") or ctx.import_id)

        for step_id, _label, error_kind, skippable in _STEPS:
            step_record = _step_record(state, step_id)
            if step_record.get("status") == "complete":
                continue
            if step_record.get("status") == "skipped" and skippable:
                continue

            state["stage"] = step_id
            _set_step_status(state, step_id, "active", started_at=time.time())
            state["progress"] = _progress_for(step_id) - 0.05
            persistence.write_state(ctx.import_id, state)

            try:
                with step_timing(ctx.import_id, step_id):
                    if step_id == "uploaded":
                        # Just record source-file metadata.
                        try:
                            stat = ctx.pptx_path.stat()
                            state["source_file"] = str(ctx.pptx_path)
                            _set_step_status(
                                state,
                                step_id,
                                "active",
                                message=f"{ctx.pptx_path.name} ({stat.st_size} bytes)",
                            )
                        except OSError:
                            pass

                    elif step_id == "analyzing":
                        manifest_data = _build_initial_manifest(ctx, render_result)
                        slide_count = len(manifest_data.get("slides") or [])
                        state["slide_count"] = slide_count
                        review["slide_count"] = slide_count

                    elif step_id == "rendering":
                        svg_dir = ctx.work_dir / "svg"
                        try:
                            render_result = renderer_pkg.render(ctx.pptx_path, svg_dir)
                        except RenderError as exc:
                            if state.get("collaboration_mode") in {"direct", "agent"}:
                                slide_count, canvas = _pptx_basic_info(ctx.pptx_path)
                                state["export_mode"] = "pptist-native"
                                state["slide_count"] = slide_count
                                review["export_mode"] = "pptist-native"  # type: ignore[typeddict-item]
                                review["slide_count"] = slide_count
                                review["canvas"] = {"width": canvas[0], "height": canvas[1]}  # type: ignore[typeddict-item]
                                _set_step_status(
                                    state,
                                    step_id,
                                    "skipped",
                                    ended_at=time.time(),
                                    message=(
                                        "Server SVG renderer unavailable; "
                                        "PPTist will parse the PPTX and provide the editable source."
                                    ),
                                )
                                state["progress"] = _progress_for(step_id)
                                state["updated_at"] = time.time()
                                state["message"] = "Awaiting PPTist client import"
                                persistence.write_state(ctx.import_id, state)
                                continue
                            raise
                        except Exception as exc:  # pragma: no cover - defensive
                            raise RenderError(
                                str(exc) or exc.__class__.__name__,
                                step_id=step_id,
                                error_kind="render",
                            ) from exc
                        state["export_mode"] = render_result.export_mode
                        state["slide_count"] = len(render_result.svg_files)
                        review["export_mode"] = render_result.export_mode  # type: ignore[typeddict-item]
                        review["slide_count"] = len(render_result.svg_files)

                    elif step_id == "detecting_assets":
                        if render_result is None:
                            # Reload from on-disk metrics + svg dir if possible.
                            render_result = _reload_render_result(ctx)
                        if render_result is None:
                            if state.get("export_mode") == "pptist-native":
                                review["assets_full"] = []  # type: ignore[typeddict-item]
                                review["assets"] = {}  # type: ignore[typeddict-item]
                                _set_step_status(
                                    state,
                                    step_id,
                                    "skipped",
                                    ended_at=time.time(),
                                    message="Skipped until the user saves a PPTist deck.",
                                )
                                state["progress"] = _progress_for(step_id)
                                state["updated_at"] = time.time()
                                persistence.write_state(ctx.import_id, state)
                                continue
                            raise ExtractionError(
                                "renderer output unavailable",
                                step_id=step_id,
                                error_kind="extraction",
                            )
                        try:
                            candidates, warnings = asset_extractor.extract_with_warnings(
                                ctx.pptx_path,
                                render_result.svg_files,
                                render_result.slide_metrics,
                            )
                        except ExtractionError:
                            raise
                        except Exception as exc:  # pragma: no cover - defensive
                            raise ExtractionError(
                                str(exc) or exc.__class__.__name__,
                                step_id=step_id,
                                error_kind="extraction",
                            ) from exc
                        warnings_aggregate.extend(warnings)
                        review["assets_full"] = [_candidate_to_full(c) for c in candidates]
                        review["assets"] = _assets_role_overrides(candidates)  # type: ignore[typeddict-item]

                        try:
                            chrome_items = chrome_detector.detect(
                                render_result.svg_files,
                                render_result.slide_metrics,
                                preserve_texts=list(review.get("preserve_texts") or []),
                                placeholder_hints={
                                    k: dict(v)
                                    for k, v in (review.get("placeholder_hints") or {}).items()
                                    if isinstance(v, dict)
                                },
                            )
                        except Exception:  # noqa: BLE001 - chrome detection is best-effort
                            chrome_items = []

                    elif step_id == "llm_review":
                        if state.get("collaboration_mode") in {"agent", "direct"} or not ctx.model_config:
                            _set_step_status(
                                state,
                                step_id,
                                "skipped",
                                ended_at=time.time(),
                                message="User/Agent workflow owns template planning.",
                            )
                            persistence.write_state(ctx.import_id, state)
                            continue
                        try:
                            client = LLMClient(ctx)
                            plan, _entry = await client.assist(
                                ctx,
                                review,
                                manifest_data,
                                feedback=None,
                                required=True,
                            )
                            merge_llm_plan(review, plan)
                        except LLMPlanError:
                            raise
                            _set_step_status(
                                state,
                                step_id,
                                "skipped",
                                ended_at=time.time(),
                                message=f"LLM unavailable: {exc.reason}",
                            )
                            persistence.write_state(ctx.import_id, state)
                            continue

                    elif step_id == "review":
                        # Build ReviewDraft + page_selections + page_type_candidates.
                        if render_result is None:
                            render_result = _reload_render_result(ctx)
                        slide_count = (
                            len(render_result.svg_files) if render_result else int(state.get("slide_count") or 0)
                        )
                        slide_metrics: list[SlideMetrics] = (
                            render_result.slide_metrics if render_result else []
                        )

                        if _mode_from_state(state) == "direct":
                            direct_order = ["cover", "toc", "chapter", "content", "ending"]
                            if slide_count == len(direct_order):
                                ordered_selections = {
                                    page_type: index + 1
                                    for index, page_type in enumerate(direct_order)
                                }
                                review["page_selections"] = ordered_selections  # type: ignore[typeddict-item]
                                review["page_type_candidates"] = {
                                    page_type: [slide_index]
                                    for page_type, slide_index in ordered_selections.items()
                                }  # type: ignore[typeddict-item]
                                review["direct_import_contract"] = {  # type: ignore[typeddict-item]
                                    "required_slide_count": len(direct_order),
                                    "order": direct_order,
                                    "templateization": "none",
                                    "design_spec_only": True,
                                }
                                review["slides"] = [  # type: ignore[typeddict-item]
                                    {
                                        "index": idx + 1,
                                        "page_type": page_type,
                                        "text_samples": [],
                                    }
                                    for idx, page_type in enumerate(direct_order)
                                ]
                            else:
                                review["page_selections"] = {}  # type: ignore[typeddict-item]
                                review["page_type_candidates"] = {pt: [] for pt in _PAGE_TYPES}  # type: ignore[typeddict-item]
                                review["direct_import_contract"] = {  # type: ignore[typeddict-item]
                                    "required_slide_count": len(direct_order),
                                    "actual_slide_count": slide_count,
                                    "order": direct_order,
                                    "templateization": "none",
                                    "design_spec_only": True,
                                    "error": "direct_import_requires_five_ordered_slides",
                                }
                                review["warnings"] = list(review.get("warnings") or []) + [  # type: ignore[typeddict-item]
                                    (
                                        "Direct import requires exactly 5 slides in order: "
                                        "cover, toc, chapter, content, ending."
                                    )
                                ]
                                review["slides"] = [  # type: ignore[typeddict-item]
                                    {
                                        "index": idx,
                                        "page_type": "content",
                                        "text_samples": [],
                                    }
                                    for idx in range(1, slide_count + 1)
                                ]

                            review["status"] = "review_required"  # type: ignore[typeddict-item]
                            persistence.write_review(ctx.import_id, review)
                            state["status"] = "review_required"
                            state["stage"] = step_id
                            state["review_required"] = True
                            state["message"] = (
                                "Direct import is ready."
                                if slide_count == len(direct_order)
                                else (
                                    "Direct import requires exactly 5 slides in order: "
                                    "cover, toc, chapter, content, ending."
                                )
                            )
                            state["progress"] = 1.0
                            state["updated_at"] = time.time()
                            if _step_record(state, step_id).get("status") == "active":
                                _set_step_status(state, step_id, "complete", ended_at=time.time())
                            persistence.write_state(ctx.import_id, state)
                            continue

                        # Run page classification when we have metrics.
                        slides_view = [
                            _slide_signal_view(m, slide_count, candidates) for m in slide_metrics
                        ]
                        try:
                            classifications = page_classifier.classify(slides_view, slide_metrics)
                            selections, _patches = page_classifier.select_representatives(
                                classifications
                            )
                        except Exception:  # noqa: BLE001 - classifier failures shouldn't block review
                            classifications = {}
                            selections = {pt: None for pt in _PAGE_TYPES}
                        if not any(selections.values()) and slide_count:
                            selections = _default_page_selections(slide_count)

                        review.setdefault("page_selections", {})
                        for pt in _PAGE_TYPES:
                            current = review["page_selections"].get(pt) if review.get("page_selections") else None  # type: ignore[union-attr]
                            if current is None and selections.get(pt) is not None:
                                review["page_selections"][pt] = selections[pt]  # type: ignore[index]

                        # Page-type candidate sets — argmax bucketing.
                        candidates_by_type: dict[str, list[int]] = {pt: [] for pt in _PAGE_TYPES}
                        for idx, cls in classifications.items():
                            if cls.page_type in candidates_by_type:
                                candidates_by_type[cls.page_type].append(int(idx))
                        for pt in _PAGE_TYPES:
                            candidates_by_type[pt].sort()
                            if not candidates_by_type[pt] and selections.get(pt):
                                candidates_by_type[pt].append(int(selections[pt]))
                        review["page_type_candidates"] = candidates_by_type  # type: ignore[typeddict-item]
                        if not review.get("slides") and slide_count:
                            by_slide = {
                                int(v): str(pt)
                                for pt, v in (review.get("page_selections") or {}).items()
                                if v
                            }
                            review["slides"] = [  # type: ignore[typeddict-item]
                                {
                                    "index": idx,
                                    "page_type": by_slide.get(idx, "content"),
                                    "text_samples": [],
                                }
                                for idx in range(1, slide_count + 1)
                            ]

                        review["status"] = "review_required"  # type: ignore[typeddict-item]
                        persistence.write_review(ctx.import_id, review)

                        state["review_required"] = True
                        state["status"] = "review_required"
                        state["template_id"] = review.get("template_id") or state.get("template_id")
                        state["progress"] = 0.9
                        state["message"] = "Awaiting review"

            except TemplateImportError as exc:
                ended = time.time()
                kind = exc.error_kind or error_kind
                _set_step_status(
                    state,
                    step_id,
                    "error",
                    ended_at=ended,
                    error=exc.reason,
                )
                state["status"] = "error"
                state["error"] = exc.reason
                state["error_kind"] = kind  # type: ignore[typeddict-item]
                state["updated_at"] = ended
                persistence.write_state(ctx.import_id, state)
                return _result_from_state(state)
            except Exception as exc:  # noqa: BLE001 - convert to TemplateImportError
                ended = time.time()
                _set_step_status(
                    state,
                    step_id,
                    "error",
                    ended_at=ended,
                    error=str(exc) or exc.__class__.__name__,
                )
                state["status"] = "error"
                state["error"] = str(exc) or exc.__class__.__name__
                state["error_kind"] = error_kind  # type: ignore[typeddict-item]
                state["updated_at"] = ended
                persistence.write_state(ctx.import_id, state)
                return _result_from_state(state)

            # Step succeeded — but only mark complete if it wasn't already
            # set to "skipped" inside the with-block.
            if _step_record(state, step_id).get("status") == "active":
                _set_step_status(state, step_id, "complete", ended_at=time.time())
            state["progress"] = _progress_for(step_id)
            state["updated_at"] = time.time()
            persistence.write_state(ctx.import_id, state)

        persistence.write_review(ctx.import_id, review)
        return _result_from_state(state)

    # ── retry_step ─────────────────────────────────────────────────────

    async def retry_step(self, import_id: str, step_id: str) -> ImportResult:
        """Reset ``step_id`` and any later steps to ``pending`` then re-run."""
        if step_id not in _STEP_IDS:
            raise ValueError(f"unknown step_id: {step_id}")
        state = persistence.read_state(import_id)
        if state is None:
            raise FileNotFoundError(import_id)

        # Validate the target step is in a retryable state.
        target = _step_record(state, step_id)
        if target.get("status") not in {"error", "skipped", "pending"}:
            # Allow retrying complete steps too — useful when a downstream
            # change requires re-running. Not currently exposed through the
            # API but keeps the function permissive.
            pass

        # Reset the target step + everything after it to pending; preserve
        # successful steps before it (and their durations).
        target_idx = _STEP_IDS.index(step_id)
        for idx, sid in enumerate(_STEP_IDS):
            if idx >= target_idx:
                record = _step_record(state, sid)
                # Preserve label, drop everything else.
                preserved_label = record.get("label", sid)
                state_steps = list(state.get("steps") or [])
                for j, r in enumerate(state_steps):
                    if isinstance(r, dict) and r.get("id") == sid:
                        state_steps[j] = {  # type: ignore[assignment]
                            "id": sid,
                            "label": preserved_label,
                            "status": "pending",
                        }
                        break
                state["steps"] = state_steps  # type: ignore[typeddict-item]

        state["status"] = "processing"
        state["error"] = None
        state["error_kind"] = None
        state["stage"] = step_id
        state["updated_at"] = time.time()
        persistence.write_state(import_id, state)

        ctx = self._ctx
        if ctx is None or ctx.import_id != import_id:
            ctx = _rebuild_context(import_id, state)
        return await self.run(ctx)

    # ── confirm ────────────────────────────────────────────────────────

    async def confirm(self, import_id: str) -> ImportResult:
        """Materialize the reviewed draft into a ``LayoutPack`` on disk."""
        state = persistence.read_state(import_id)
        review = persistence.read_review(import_id)
        if state is None or review is None:
            raise FileNotFoundError(import_id)

        template_id = str(review.get("template_id") or state.get("template_id") or "")
        if not template_id:
            raise PersistenceError(
                "missing template_id",
                step_id="confirm",
                error_kind="persistence",
            )

        # Idempotency short-circuit: a complete state with the layout dir
        # already on disk → return as-is without rewriting.
        if state.get("status") == "complete":
            existing = persistence.load_template_pack(template_id)
            if existing is not None:
                return _result_from_state(state)

        ctx = self._ctx
        if ctx is None or ctx.import_id != import_id:
            ctx = _rebuild_context(import_id, state)

        mode = _mode_from_state(state)
        if mode == "direct":
            return _confirm_direct_ordered_template(ctx, state, review, template_id)
        if mode == "agent":
            return _confirm_agent_template(ctx, state, review, template_id)

        # Reload renderer outputs from disk.
        render_result = _reload_render_result(ctx)
        if render_result is None:
            deck = _load_pptist_deck(ctx)
            if deck is not None:
                return _confirm_from_pptist_deck(ctx, state, review, template_id, deck)
            raise PersistenceError(
                "rendered SVGs unavailable; save the PPTist deck or install a server-side PPTX renderer",
                step_id="confirm",
                error_kind="persistence",
            )

        slide_metrics_by_index: dict[int, SlideMetrics] = {
            m.slide_index: m for m in render_result.slide_metrics
        }

        # Build asset role lookups (file_name → role).
        assets_full = list(review.get("assets_full") or [])
        asset_overrides = dict(review.get("assets") or {})
        asset_role_by_file: dict[str, str] = {}
        common_assets: list[ManifestAssetEntry] = []
        layout_assets: dict[str, bytes] = {}
        canvas_w = render_result.canvas[0] or 1280
        canvas_h = render_result.canvas[1] or 720

        for asset in assets_full:
            asset_id = asset.get("asset_id")
            file_name = asset.get("file_name") or ""
            override = asset_overrides.get(asset_id) if asset_id else None
            role = (
                (override or {}).get("role")
                if isinstance(override, dict)
                else None
            ) or asset.get("recommended_role") or "decoration"
            asset_role_by_file[file_name] = role
            if role in ("logo", "background", "decoration"):
                # First occurrence drives the geometry.
                first_occ = (asset.get("occurrences") or [{}])[0]
                bbox = {
                    "x": float(first_occ.get("x") or 0.0) / max(1.0, float(canvas_w)),
                    "y": float(first_occ.get("y") or 0.0) / max(1.0, float(canvas_h)),
                    "width": float(first_occ.get("width") or 0.0) / max(1.0, float(canvas_w)),
                    "height": float(first_occ.get("height") or 0.0) / max(1.0, float(canvas_h)),
                }
                entry: ManifestAssetEntry = {  # type: ignore[typeddict-unknown-key]
                    "asset_id": str(asset_id or ""),
                    "file_name": str(file_name),
                    "role": role,  # type: ignore[typeddict-item]
                    "pages": list(asset.get("pages") or []),
                    "bbox_norm": bbox,
                    "sha1": str(asset.get("sha1") or ""),
                    "bytes": int(asset.get("bytes") or 0),
                }
                common_assets.append(entry)
                # Copy bytes from work_dir/assets if present.
                src_path = ctx.work_dir / "assets" / file_name
                if src_path.exists():
                    try:
                        layout_assets[file_name] = src_path.read_bytes()
                    except OSError:
                        pass

        # Element actions filtered by page_type.
        all_actions: list[ElementActionRecord] = list(review.get("element_actions") or [])

        page_selections = dict(review.get("page_selections") or {})
        templateize_outputs: dict[str, TemplateizeOutput] = {}
        warnings_aggregate: list[ManifestWarning] = []
        content_area_box: BoundingBox | None = None
        export_mode = str(review.get("export_mode") or render_result.export_mode)
        svg_dir = ctx.work_dir / "svg"

        for page_type in _PAGE_TYPES:
            slide_index = page_selections.get(page_type)
            if not slide_index:
                continue
            svg_path = svg_dir / f"slide-{int(slide_index):03d}.svg"
            if not svg_path.exists():
                continue
            metric = slide_metrics_by_index.get(int(slide_index))
            if metric is None:
                metric = SlideMetrics(
                    slide_index=int(slide_index),
                    canvas_width=canvas_w,
                    canvas_height=canvas_h,
                    font_substitutions={},
                    rendered_at=time.time(),
                )

            actions: list[ElementAction] = []
            for record in all_actions:
                if not isinstance(record, dict):
                    continue
                if record.get("page_type") != page_type:
                    continue
                actions.append(
                    ElementAction(
                        page_type=record.get("page_type", page_type),  # type: ignore[arg-type]
                        element_id=str(record.get("element_id") or ""),
                        action=record.get("action", "keep"),  # type: ignore[arg-type]
                        placeholder=record.get("placeholder"),
                        reason=record.get("reason"),
                        source=record.get("source", "rule"),  # type: ignore[arg-type]
                    )
                )

            # Run chrome detector on the fly — cheap and avoids stale cache.
            try:
                chrome_items = chrome_detector.detect(
                    [svg_path],
                    [metric],
                    preserve_texts=list(review.get("preserve_texts") or []),
                    placeholder_hints={
                        k: dict(v)
                        for k, v in (review.get("placeholder_hints") or {}).items()
                        if isinstance(v, dict)
                    },
                )
            except Exception:  # noqa: BLE001
                chrome_items = []

            try:
                output = templateizer.templateize(
                    svg_path,
                    page_type,  # type: ignore[arg-type]
                    metrics=metric,
                    actions=actions,
                    chrome_items=chrome_items,
                    asset_roles=asset_role_by_file,
                )
            except Exception as exc:  # noqa: BLE001
                raise PersistenceError(
                    f"templateize failed for {page_type}: {exc}",
                    step_id="confirm",
                    error_kind="persistence",
                ) from exc

            templateize_outputs[page_type] = output
            warnings_aggregate.extend(output.warnings or [])
            if page_type == "content" and output.content_area is not None:
                content_area_box = output.content_area

        # Inline image references in templated SVGs so the persisted pack
        # never points at ``http(s)://`` URLs and renders identically in
        # the preview and in production (Requirement 12.2; Property 22).
        for page_type, output in templateize_outputs.items():
            output.svg_text = inline_svg_asset_refs(output.svg_text, layout_assets)

        # Assemble manifest.
        canvas_block: ManifestCanvas = {"width": int(canvas_w), "height": int(canvas_h)}
        theme_block: ManifestTheme = {  # type: ignore[typeddict-unknown-key]
            "colors": {},
            "fonts": {},
        }
        # Pull theme colors if previously captured on the state record.
        theme_colors = list(state.get("theme_colors") or [])
        for idx, color in enumerate(theme_colors[:6]):
            theme_block["colors"][f"accent{idx + 1}"] = str(color)

        content_area_dict: ManifestContentArea | None = None
        if content_area_box is not None:
            content_area_dict = {
                "x": float(content_area_box.x),
                "y": float(content_area_box.y),
                "width": float(content_area_box.width),
                "height": float(content_area_box.height),
            }

        page_selections_clean = {
            pt: int(v) if v is not None else None for pt, v in page_selections.items() if pt in _PAGE_TYPES
        }
        page_type_candidates = {
            pt: list(v) for pt, v in (review.get("page_type_candidates") or {}).items() if pt in _PAGE_TYPES
        }

        manifest: TemplateManifest = {  # type: ignore[typeddict-unknown-key]
            "schema_version": 2,
            "template_id": template_id,
            "label": str(review.get("label") or state.get("label") or template_id),
            "source_file": str(state.get("source_file") or ctx.pptx_path.name),
            "slide_count": int(state.get("slide_count") or render_result and len(render_result.svg_files) or 0),
            "canvas": canvas_block,
            "theme": theme_block,
            "page_type_candidates": page_type_candidates,  # type: ignore[typeddict-item]
            "page_selections": page_selections_clean,  # type: ignore[typeddict-item]
            "common_assets": common_assets,
            "content_area": content_area_dict,
            "warnings": warnings_aggregate,
            "imported_at": time.time(),
            "export_mode": export_mode,  # type: ignore[typeddict-item]
            "importer_version": _IMPORTER_VERSION,
            "source_kind": "pptist_import",
            "import_mode": str(state.get("collaboration_mode") or "direct"),
            "clean_generation_ready": True,
        }

        svgs: dict[str, str] = {pt: out.svg_text for pt, out in templateize_outputs.items()}
        if str(state.get("collaboration_mode") or "") == "agent":
            agent_svgs = _load_agent_template_svgs(ctx)
            if agent_svgs:
                svgs.update(agent_svgs)

        design_spec = str(review.get("design_spec_md") or "").strip()
        if not design_spec and str(state.get("collaboration_mode") or "") == "agent":
            design_spec = _load_agent_design_spec(ctx)
        if not design_spec:
            design_spec = _default_design_spec(review, manifest)

        import_trace: list[LLMTraceEntry] = list(review.get("llm_trace") or [])

        pack = LayoutPack(
            template_id=template_id,
            label=manifest["label"],
            manifest=manifest,
            svgs=svgs,  # type: ignore[arg-type]
            design_spec=design_spec,
            assets=layout_assets,
            import_trace=import_trace,
        )

        try:
            install_result = persistence.transactional_install_layout(template_id, pack)
        except PersistenceError:
            raise
        except Exception as exc:  # noqa: BLE001
            raise PersistenceError(
                str(exc) or exc.__class__.__name__,
                step_id="confirm",
                error_kind="persistence",
            ) from exc

        try:
            persistence.update_user_index(template_id, manifest["label"], manifest)
        except Exception as exc:  # noqa: BLE001
            raise PersistenceError(
                f"index update failed: {exc}",
                step_id="confirm",
                error_kind="persistence",
            ) from exc

        # Finalize state.
        state["status"] = "complete"
        state["stage"] = "complete"
        state["progress"] = 1.0
        state["template_id"] = template_id
        state["label"] = manifest["label"]
        state["export_mode"] = export_mode
        state["review_required"] = False
        state["error"] = None
        state["error_kind"] = None
        state["updated_at"] = time.time()
        persistence.write_state(import_id, state)

        result = _result_from_state(state)
        if install_result.already_complete:
            result.message = "Template already installed."
        return result


# ─────────────────────────────────────────────────────────────────────────────
# Module-level helpers
# ─────────────────────────────────────────────────────────────────────────────


def _rebuild_context(import_id: str, state: ImportTaskState) -> PipelineContext:
    """Reconstruct a :class:`PipelineContext` from on-disk state.

    Used when ``retry_step`` / ``confirm`` is called without a live
    pipeline reference (e.g. fresh worker process).
    """
    from backend.config import settings

    work_dir = settings.workspaces_dir / "template_imports" / import_id
    pptx_path = Path(state.get("source_file") or "")
    return PipelineContext(
        import_id=import_id,
        work_dir=work_dir,
        pptx_path=pptx_path,
        label=str(state.get("label") or import_id),
        model_config={},
    )


def _reload_render_result(ctx: PipelineContext) -> RenderResult | None:
    """Best-effort reload of the rendered SVG batch from disk."""
    from .renderer.slide_metrics import read_slide_metrics

    svg_dir = ctx.work_dir / "svg"
    if not svg_dir.is_dir():
        return None
    svg_files = sorted(svg_dir.glob("slide-*.svg"))
    if not svg_files:
        return None
    metrics = read_slide_metrics(svg_dir)
    if not metrics:
        # Synthesize minimal metrics from filename order if metadata is missing.
        canvas = (1280, 720)
        metrics = [
            SlideMetrics(
                slide_index=idx + 1,
                canvas_width=canvas[0],
                canvas_height=canvas[1],
                font_substitutions={},
                rendered_at=0.0,
            )
            for idx in range(len(svg_files))
        ]
    canvas_w = metrics[0].canvas_width if metrics else 1280
    canvas_h = metrics[0].canvas_height if metrics else 720
    return RenderResult(
        svg_files=svg_files,
        slide_metrics=metrics,
        export_mode="libreoffice",  # best-guess; actual export_mode is recorded on state
        canvas=(canvas_w, canvas_h),
    )


def _default_design_spec(review: ReviewDraft, manifest: TemplateManifest) -> str:
    """Synthesize a minimal ``design_spec.md`` when none was provided."""
    lines = [
        f"# {manifest.get('label') or manifest.get('template_id')}",
        "",
        f"- Imported from: {manifest.get('source_file') or ''}",
        f"- Slide count: {manifest.get('slide_count', 0)}",
        f"- Export mode: {manifest.get('export_mode') or ''}",
        "",
        "## Page Selections",
    ]
    for pt in _PAGE_TYPES:
        sel = (manifest.get("page_selections") or {}).get(pt)
        lines.append(f"- {pt}: {sel if sel else '—'}")
    return "\n".join(lines) + "\n"


def preview_templateized(import_id: str) -> dict[str, Any]:
    """Build a clean, non-installed template preview for the current review."""
    state = persistence.read_state(import_id)
    review = persistence.read_review(import_id)
    if state is None or review is None:
        raise FileNotFoundError(import_id)
    ctx = _rebuild_context(import_id, state)
    mode = _mode_from_state(state)
    if mode == "agent":
        return _preview_agent_template(ctx, state, review)

    render_result = _reload_render_result(ctx)
    if mode == "direct":
        deck = _load_pptist_deck(ctx)
        if render_result is not None:
            return _direct_preview_from_rendered(ctx, state, review, render_result)
        if deck is not None:
            return _direct_preview_from_pptist_deck(ctx, state, review, deck)
        raise PersistenceError(
            "Direct import requires exactly 5 saved PPTist slides, or rendered SVGs from a 5-slide PPTX.",
            step_id="preview",
            error_kind="persistence",
        )

    if render_result is None:
        deck = _load_pptist_deck(ctx)
        if deck is not None:
            return _preview_from_pptist_deck(ctx, state, review, deck)
        raise PersistenceError(
            "rendered SVGs unavailable",
            step_id="preview",
            error_kind="persistence",
        )

    canvas_w = render_result.canvas[0] or 1280
    canvas_h = render_result.canvas[1] or 720
    slide_metrics_by_index = {m.slide_index: m for m in render_result.slide_metrics}
    assets_full = list(review.get("assets_full") or [])
    asset_overrides = dict(review.get("assets") or {})
    asset_role_by_file: dict[str, str] = {}
    layout_assets: dict[str, bytes] = {}
    for asset in assets_full:
        asset_id = asset.get("asset_id")
        file_name = str(asset.get("file_name") or "")
        override = asset_overrides.get(asset_id) if asset_id else None
        role = (
            (override or {}).get("role")
            if isinstance(override, dict)
            else None
        ) or asset.get("recommended_role") or "decoration"
        asset_role_by_file[file_name] = str(role)
        if role in ("logo", "background", "decoration"):
            src_path = ctx.work_dir / "assets" / file_name
            if src_path.exists():
                try:
                    layout_assets[file_name] = src_path.read_bytes()
                except OSError:
                    pass

    all_actions: list[ElementActionRecord] = list(review.get("element_actions") or [])
    page_selections = dict(review.get("page_selections") or {})
    svg_dir = ctx.work_dir / "svg"
    out: dict[str, Any] = {
        "template_id": str(review.get("template_id") or state.get("template_id") or import_id),
        "label": str(review.get("label") or state.get("label") or import_id),
        "cover_svg": "",
        "toc_svg": "",
        "chapter_svg": "",
        "content_svg": "",
        "ending_svg": "",
        "design_spec": str(review.get("design_spec_md") or ""),
        "theme_colors": list(state.get("theme_colors") or []),
        "warnings": [],
        "placeholders": {},
    }

    for page_type in _PAGE_TYPES:
        slide_index = page_selections.get(page_type)
        if not slide_index:
            continue
        svg_path = svg_dir / f"slide-{int(slide_index):03d}.svg"
        if not svg_path.exists():
            continue
        metric = slide_metrics_by_index.get(int(slide_index)) or SlideMetrics(
            slide_index=int(slide_index),
            canvas_width=canvas_w,
            canvas_height=canvas_h,
            font_substitutions={},
            rendered_at=time.time(),
        )
        actions: list[ElementAction] = []
        for record in all_actions:
            if not isinstance(record, dict) or record.get("page_type") != page_type:
                continue
            actions.append(
                ElementAction(
                    page_type=record.get("page_type", page_type),  # type: ignore[arg-type]
                    element_id=str(record.get("element_id") or ""),
                    action=record.get("action", "keep"),  # type: ignore[arg-type]
                    placeholder=record.get("placeholder"),
                    reason=record.get("reason"),
                    source=record.get("source", "rule"),  # type: ignore[arg-type]
                )
            )
        try:
            chrome_items = chrome_detector.detect(
                [svg_path],
                [metric],
                preserve_texts=list(review.get("preserve_texts") or []),
                placeholder_hints={
                    k: dict(v)
                    for k, v in (review.get("placeholder_hints") or {}).items()
                    if isinstance(v, dict)
                },
            )
        except Exception:
            chrome_items = []
        output = templateizer.templateize(
            svg_path,
            page_type,  # type: ignore[arg-type]
            metrics=metric,
            actions=actions,
            chrome_items=chrome_items,
            asset_roles=asset_role_by_file,
        )
        svg_text = inline_svg_asset_refs(output.svg_text, layout_assets)
        out[f"{page_type}_svg"] = svg_text
        out["warnings"].extend([dict(item) for item in output.warnings or []])
        out["placeholders"][page_type] = [
            {
                "name": ph.name,
                "bbox": {
                    "x": ph.bbox.x,
                    "y": ph.bbox.y,
                    "width": ph.bbox.width,
                    "height": ph.bbox.height,
                },
                "style": ph.style,
            }
            for ph in output.placeholders
        ]
    if str(state.get("collaboration_mode") or "") == "agent":
        agent_svgs = _load_agent_template_svgs(ctx)
        if agent_svgs:
            for page_type, svg_text in agent_svgs.items():
                out[f"{page_type}_svg"] = svg_text
            agent_spec = _load_agent_design_spec(ctx)
            if agent_spec:
                out["design_spec"] = agent_spec
            out["warnings"].append(
                {
                    "code": "agent_template_preview",
                    "message": "Preview includes Agent-authored template SVG outputs.",
                }
            )
    return out


__all__ = ["Pipeline", "inline_svg_asset_refs", "preview_templateized"]
